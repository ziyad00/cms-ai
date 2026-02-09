import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


def main(argv: list[str]) -> int:
    # Accept either 2 or 4 arguments (with optional --company-info)
    if len(argv) not in [2, 4]:
        print("Usage: render_pptx.py <spec.json> <out.pptx> [--company-info <company.json>]", file=sys.stderr)
        return 2

    spec_path = Path(argv[0])
    out_path = Path(argv[1])

    # Optionally handle --company-info argument
    company_info = None
    if len(argv) == 4 and argv[2] == "--company-info":
        try:
            company_info_path = Path(argv[3])
            company_info = json.loads(company_info_path.read_text(encoding="utf-8"))
        except Exception:
            pass  # Ignore company info parsing errors

    spec = json.loads(spec_path.read_text(encoding="utf-8"))
    prs = Presentation()

    # Deterministic deck: one slide per layout.
    # We map placeholders roughly to text boxes.
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    for layout in spec.get("layouts", []):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        title = layout.get("name", "Layout")

        # Title banner
        tb = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.2), slide_w - Inches(1.0), Inches(0.5)
        )
        tf = tb.text_frame
        tf.text = title
        tf.paragraphs[0].font.size = Pt(20)

        for ph in layout.get("placeholders", []):
            geom = ph.get("geometry") or {}
            x = float(geom.get("x", 0.1))
            y = float(geom.get("y", 0.2))
            w = float(geom.get("w", 0.8))
            h = float(geom.get("h", 0.2))

            left = int(slide_w * x)
            top = int(slide_h * y)
            width = int(slide_w * w)
            height = int(slide_h * h)

            if ph.get("type") == "text":
                box = slide.shapes.add_textbox(left, top, width, height)
                # Use actual content if available, otherwise fall back to placeholder ID
                text = ph.get("content") or f"[{ph.get('id', 'placeholder')}]"
                box.text_frame.text = text

                # Basic text formatting
                if ph.get("id") in ["title", "heading"]:
                    box.text_frame.paragraphs[0].font.size = Pt(24)
                    box.text_frame.paragraphs[0].font.bold = True
                elif ph.get("id") in ["subtitle", "subheading"]:
                    box.text_frame.paragraphs[0].font.size = Pt(18)
                else:
                    box.text_frame.paragraphs[0].font.size = Pt(14)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
