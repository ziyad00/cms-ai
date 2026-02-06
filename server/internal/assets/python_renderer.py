#!/usr/bin/env python3
"""
Simple Python PPTX Renderer - embedded version
"""

import sys
import json
import argparse
import asyncio
import os
import logging
from pathlib import Path
from typing import Dict, Any, Optional

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError as e:
    print(f"ERROR: python-pptx library is required. Install with: pip install python-pptx. Error: {e}", file=sys.stderr)
    sys.exit(1)

logging.basicConfig(level=logging.INFO)

class SimplePPTXRenderer:
    def __init__(self):
        self.logger = logging.getLogger(__name__)

    def render_pptx(self, spec_data, output_path):
        """Simple PPTX rendering without AI features"""
        prs = Presentation()

        # Remove default slide if any
        if len(prs.slides) > 0:
            slide_to_remove = prs.slides[0]
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]

        # Process layouts
        layouts = spec_data.get('layouts', [])
        for layout in layouts:
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

            # Add content
            placeholders = layout.get('placeholders', [])
            for ph in placeholders:
                content = ph.get('content', '')

                # Use geometry if available
                if 'geometry' in ph:
                    x = ph['geometry'].get('x', 1.0)
                    y = ph['geometry'].get('y', 1.0)
                    w = ph['geometry'].get('w', 8.0)
                    h = ph['geometry'].get('h', 1.0)
                else:
                    x = ph.get('x', 1.0)
                    y = ph.get('y', 1.0)
                    w = ph.get('width', 8.0)
                    h = ph.get('height', 1.0)

                # Add text box
                text_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
                text_frame = text_box.text_frame
                text_frame.clear()

                p = text_frame.paragraphs[0]
                p.text = content

        # Save presentation
        prs.save(output_path)
        self.logger.info(f"Generated: {output_path}")

async def main():
    parser = argparse.ArgumentParser(description='Embedded PPTX Renderer v2', prog='embedded_renderer')
    parser.add_argument('spec_file', help='JSON spec file')
    parser.add_argument('output_file', help='Output PPTX file')
    parser.add_argument('--company-info', help='Company info JSON file (optional)')
    parser.add_argument('--hf-api-key', help='Hugging Face API key (or set HUGGING_FACE_API_KEY env var)')

    args = parser.parse_args()

    try:
        # Load spec
        with open(args.spec_file, 'r') as f:
            spec_data = json.load(f)

        # Render with simple renderer
        renderer = SimplePPTXRenderer()
        renderer.render_pptx(spec_data, args.output_file)

        print(f"✅ Generated: {args.output_file}")

    except Exception as e:
        print(f"❌ Error: {e}", file=sys.stderr)
        sys.exit(1)

def main_sync():
    """Synchronous entry point for backward compatibility"""
    asyncio.run(main())

if __name__ == '__main__':
    main_sync()