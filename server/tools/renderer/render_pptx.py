#!/usr/bin/env python3
"""
AI-Powered Python PPTX Renderer for CMS-AI v2
Integrates olama's visual rendering with Hugging Face AI design analysis
Generates presentations with intelligent design decisions based on content
"""

import re
import sys
import json
import argparse
import asyncio
import os
import logging
from pathlib import Path
from typing import Dict, Any, Optional, List, Tuple

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
except ImportError as e:
    print(f"ERROR: python-pptx library is required. Install with: pip install python-pptx. Error: {e}", file=sys.stderr)
    sys.exit(1)

# Import olama's AI and design modules (local copies)
try:
    from ai_design_generator import AIDesignGenerator
    from design_templates import DesignTemplateLibrary, get_design_system_for_content
    from abstract_background_renderer import CompositeBackgroundRenderer
except ImportError as e:
    print(f"ERROR: Failed to import olama modules: {e}", file=sys.stderr)
    sys.exit(1)
except Exception as e:
    print(f"ERROR: Unexpected error during imports: {e}", file=sys.stderr)
    sys.exit(1)

logging.basicConfig(level=logging.INFO)

# Standard slide dimensions in inches
SLIDE_WIDTH_INCHES = 10.0
SLIDE_HEIGHT_INCHES = 7.5


class ContentType:
    """Types of content detected by analysis (matches olama ContentType)."""
    TEXT_HEAVY = "text_heavy"
    DATA_DRIVEN = "data_driven"
    LIST_ITEMS = "list_items"
    COMPARISON = "comparison"
    TIMELINE = "timeline"
    HIERARCHY = "hierarchy"
    QUOTE = "quote"
    IMAGE_TEXT = "image_text"


class ContentAnalyzer:
    """Analyzes slide content to determine optimal design approach (ported from olama)."""

    def analyze_content(self, title: str, content_items: List[str]) -> Dict[str, Any]:
        """Analyze slide content and return analysis dict."""
        content_text = ' '.join([title] + [str(item) for item in content_items])
        word_count = len(content_text.split())

        content_type = self._detect_content_type(title, content_items)
        complexity = self._analyze_complexity(content_items)
        sentiment = self._analyze_sentiment(content_text)
        key_concepts = self._extract_key_concepts(content_text)
        has_numbers = self._has_numbers(content_text)
        has_dates = self._has_dates(content_text)
        hierarchy_level = self._determine_hierarchy_level(title)
        visual_weight = self._calculate_visual_weight(word_count, complexity, hierarchy_level)

        return {
            "content_type": content_type,
            "word_count": word_count,
            "complexity": complexity,
            "sentiment": sentiment,
            "key_concepts": key_concepts,
            "has_numbers": has_numbers,
            "has_dates": has_dates,
            "hierarchy_level": hierarchy_level,
            "visual_weight": visual_weight,
        }

    def _detect_content_type(self, title: str, content: List[str]) -> str:
        """Detect what type of content this slide contains."""
        title_lower = title.lower()
        content_text = ' '.join([str(item) for item in content]).lower()

        if any(word in title_lower for word in ['timeline', 'phases', 'roadmap', 'schedule']):
            return ContentType.TIMELINE
        if any(word in title_lower for word in ['metrics', 'kpi', 'results', 'analysis']):
            return ContentType.DATA_DRIVEN
        if any(word in title_lower for word in ['vs', 'comparison', 'options', 'alternatives']):
            return ContentType.COMPARISON
        if any(word in title_lower for word in ['architecture', 'structure', 'hierarchy', 'organization']):
            return ContentType.HIERARCHY
        if len(content) == 1 and len(content_text.split()) < 30:
            return ContentType.QUOTE
        if len(content) > 3:
            return ContentType.LIST_ITEMS
        if len(content_text.split()) > 100:
            return ContentType.TEXT_HEAVY
        return ContentType.LIST_ITEMS

    def _analyze_complexity(self, content: List[str]) -> str:
        """Analyze content complexity based on average word length per item."""
        if not content:
            return "simple"
        total_words = sum(len(str(item).split()) for item in content)
        avg_length = total_words / len(content)
        if avg_length > 15:
            return "complex"
        elif avg_length > 8:
            return "medium"
        return "simple"

    def _analyze_sentiment(self, text: str) -> str:
        """Analyze emotional tone of content."""
        text_lower = text.lower()
        urgent_words = ['critical', 'urgent', 'immediate', 'risk', 'threat', 'warning']
        if any(word in text_lower for word in urgent_words):
            return "urgent"
        positive_words = ['success', 'growth', 'achievement', 'improvement', 'benefit']
        positive_count = sum(1 for word in positive_words if word in text_lower)
        negative_words = ['problem', 'issue', 'challenge', 'failure', 'decline']
        negative_count = sum(1 for word in negative_words if word in text_lower)
        if positive_count > negative_count:
            return "positive"
        elif negative_count > positive_count:
            return "negative"
        return "neutral"

    def _extract_key_concepts(self, text: str) -> List[str]:
        """Extract key concepts (capitalized words + acronyms)."""
        words = re.findall(r'\b[A-Z][a-z]+\b', text)
        technical_terms = re.findall(r'\b[A-Z]{2,}\b', text)
        return list(set(words + technical_terms))[:5]

    def _has_numbers(self, text: str) -> bool:
        return bool(re.search(r'\d+[%$]?|\d+\.\d+', text))

    def _has_dates(self, text: str) -> bool:
        date_patterns = [
            r'\d{4}[-/]\d{1,2}[-/]\d{1,2}',
            r'\d{1,2}[-/]\d{1,2}[-/]\d{4}',
            r'(Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)',
            r'Q[1-4]',
        ]
        return any(re.search(pattern, text, re.IGNORECASE) for pattern in date_patterns)

    def _determine_hierarchy_level(self, title: str) -> int:
        """1=executive, 2=section, 3=detail."""
        title_lower = title.lower()
        if any(word in title_lower for word in ['executive', 'summary', 'overview']):
            return 1
        if any(word in title_lower for word in ['introduction', 'conclusion', 'agenda']):
            return 2
        return 3

    def _calculate_visual_weight(self, word_count: int, complexity: str, hierarchy: int) -> float:
        """Calculate visual emphasis needed (0.1-1.0)."""
        weight = 0.5
        if word_count < 20:
            weight += 0.3
        elif word_count > 80:
            weight -= 0.2
        complexity_weights = {"simple": 0.2, "medium": 0.0, "complex": -0.2}
        weight += complexity_weights.get(complexity, 0)
        hierarchy_weights = {1: 0.4, 2: 0.2, 3: 0.0}
        weight += hierarchy_weights.get(hierarchy, 0)
        return max(0.1, min(1.0, weight))


class SmartDesignRules:
    """Professional design rules engine (ported from olama)."""

    @staticmethod
    def calculate_golden_ratio_spacing(base_size: float) -> Dict[str, float]:
        """Calculate spacing based on golden ratio (1.618)."""
        phi = 1.618
        return {
            'tight': base_size / phi,
            'normal': base_size,
            'loose': base_size * phi,
            'extra_loose': base_size * phi * phi,
        }

    @staticmethod
    def calculate_optimal_font_sizes(
        visual_weight: float, word_count: int, content_type: str
    ) -> Dict[str, int]:
        """Calculate optimal font sizes based on content analysis."""
        base_title = 32
        base_body = 18

        title_size = int(base_title * (0.7 + visual_weight * 0.8))
        body_size = int(base_body * (0.8 + visual_weight * 0.4))

        if word_count < 20:
            title_size += 8
            body_size += 4
        elif word_count > 80:
            title_size -= 4
            body_size -= 2

        if content_type == ContentType.QUOTE:
            body_size += 6
        elif content_type == ContentType.DATA_DRIVEN:
            body_size -= 2

        return {
            'title': max(24, min(48, title_size)),
            'body': max(12, min(24, body_size)),
            'caption': max(10, min(14, body_size - 4)),
        }

    @staticmethod
    def ensure_contrast_ratio(background_color: str, text_color: str, min_ratio: float = 4.5) -> str:
        """Ensure sufficient contrast for readability (WCAG simplified)."""
        if background_color.upper().startswith('#F') or background_color.upper() == '#FFFFFF':
            return '#1F1F1F'
        return '#FFFFFF'

    @staticmethod
    def calculate_dynamic_positioning(
        content_items: List[str], slide_width: float, slide_height: float,
        has_title: bool = True
    ) -> Dict[str, float]:
        """Calculate dynamic positioning based on content volume (ported from olama).

        Returns position dict with x, y, width, height for the body content area.
        """
        item_count = len(content_items)
        total_chars = sum(len(str(item)) for item in content_items)

        # Title area takes top 20%
        title_height = slide_height * 0.20 if has_title else 0
        margin = 0.5

        # Body area
        body_y = title_height + margin * 0.5 if has_title else margin
        body_height = slide_height - body_y - margin

        # Width adjusts based on content density
        if total_chars > 500 or item_count > 8:
            # Dense content: use full width
            body_x = margin
            body_width = slide_width - 2 * margin
        elif item_count <= 1:
            # Sparse content: center with generous margins
            body_x = slide_width * 0.15
            body_width = slide_width * 0.7
        else:
            # Normal content
            body_x = margin + 0.3
            body_width = slide_width - 2 * margin - 0.6

        return {
            'x': body_x,
            'y': body_y,
            'width': body_width,
            'height': body_height,
        }


class AdvancedLayoutEngine:
    """Advanced layout intelligence for optimal content arrangement (ported from olama)."""

    @staticmethod
    def detect_layout_pattern(content_items: List[str]) -> str:
        """Detect optimal layout pattern based on content structure."""
        if not content_items:
            return "single_column"

        item_count = len(content_items)
        joined = " ".join(content_items).lower()

        has_vs = any("vs" in item.lower() for item in content_items)
        has_versus = any(w in joined for w in ["versus", "compared to", "current", "proposed"])
        has_numbers = any(c.isdigit() for item in content_items for c in item)
        has_pct = any("%" in item for item in content_items)
        has_dates = any(w in joined for w in ["q1", "q2", "q3", "q4", "2024", "2025", "2026"])
        has_timeline = any(w in joined for w in ["phase", "milestone", "timeline", "schedule"])

        if has_vs or has_versus:
            return "comparison_columns"
        elif has_dates or has_timeline:
            return "horizontal_timeline"
        elif has_numbers and has_pct and item_count <= 6:
            return "metrics_grid"
        elif item_count == 2:
            return "two_column"
        elif item_count == 3:
            return "three_column"
        elif 4 <= item_count <= 6:
            return "grid_layout"
        elif item_count > 8:
            return "multi_column"
        return "single_column"

    @staticmethod
    def calculate_optimal_columns(item_count: int, slide_width: float) -> Dict[str, Any]:
        """Calculate optimal column layout."""
        if item_count <= 2:
            columns, column_width = 1, slide_width * 0.8
        elif item_count <= 4:
            columns, column_width = 2, slide_width * 0.4
        elif item_count <= 9:
            columns, column_width = 3, slide_width * 0.25
        else:
            columns, column_width = 4, slide_width * 0.2
        return {
            "columns": columns,
            "column_width": column_width,
            "spacing": slide_width * 0.02,
            "margin": slide_width * 0.05,
        }


class DynamicChartGenerator:
    """Auto-generate charts from data content (ported from olama smart_slide_generator)."""

    @staticmethod
    def detect_data_pattern(content_items: List[str]) -> Dict[str, Any]:
        """Detect data patterns that can be visualized as charts."""
        data_patterns = {
            "chart_type": None,
            "data": [],
            "labels": [],
            "has_percentages": False,
            "has_timeline": False,
        }

        percentage_pattern = r'(\d+(?:\.\d+)?)\s*%'
        numeric_pattern = r'(\d+(?:,\d{3})*(?:\.\d+)?)'

        for item in content_items:
            # Extract percentages
            percentages = re.findall(percentage_pattern, item)
            if percentages:
                data_patterns["has_percentages"] = True
                for pct in percentages:
                    data_patterns["data"].append(float(pct))
                    label = re.sub(percentage_pattern, '', item).strip().rstrip(':')
                    data_patterns["labels"].append(label[:30])  # Truncate long labels

            # Check for timeline data
            timeline_matches = re.findall(r'(Q[1-4]|[A-Za-z]+)\s*\d{4}', item)
            if timeline_matches:
                data_patterns["has_timeline"] = True

            # Extract general numeric data (skip if already counted as percentage)
            if not percentages:
                numbers = re.findall(numeric_pattern, item)
                for num in numbers:
                    try:
                        value = float(num.replace(',', ''))
                        data_patterns["data"].append(value)
                        label = re.sub(numeric_pattern, '', item).strip().rstrip(':')
                        data_patterns["labels"].append(label[:30])
                    except ValueError:
                        continue

        # Determine chart type
        if data_patterns["has_percentages"] and len(data_patterns["data"]) >= 2:
            data_patterns["chart_type"] = "pie"
        elif data_patterns["has_timeline"] and len(data_patterns["data"]) >= 2:
            data_patterns["chart_type"] = "line"
        elif len(data_patterns["data"]) >= 2:
            data_patterns["chart_type"] = "bar"

        return data_patterns

    @staticmethod
    def create_chart(slide, chart_data: Dict[str, Any], x, y, width, height):
        """Create a chart from detected data patterns."""
        if not chart_data["data"] or not chart_data["labels"]:
            return None

        try:
            data_obj = CategoryChartData()
            data_obj.categories = chart_data["labels"][:len(chart_data["data"])]
            data_obj.add_series('Values', chart_data["data"])

            if chart_data["chart_type"] == "pie":
                chart_type = XL_CHART_TYPE.PIE
            elif chart_data["chart_type"] == "line":
                chart_type = XL_CHART_TYPE.LINE
            else:
                chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED

            chart = slide.shapes.add_chart(
                chart_type, x, y, width, height, data_obj
            ).chart
            return chart
        except Exception as e:
            logging.getLogger(__name__).warning(f"Could not create chart: {e}")
            return None

    @staticmethod
    def create_progress_bar(slide, x, y, width, height, percentage: float,
                            label: str, colors: Dict[str, str]):
        """Create a visual progress bar for percentage values."""
        # Background bar
        bg_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, width, height)
        bg_bar.fill.solid()
        light_color = colors.get('light', '#E5E5E5').lstrip('#')
        bg_bar.fill.fore_color.rgb = RGBColor.from_string(light_color)
        bg_bar.line.fill.background()

        # Filled portion
        filled_width = int(width * (percentage / 100))
        if filled_width > 0:
            progress_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x, y, filled_width, height
            )
            progress_bar.fill.solid()
            accent_color = colors.get('accent', '#00B050').lstrip('#')
            progress_bar.fill.fore_color.rgb = RGBColor.from_string(accent_color)
            progress_bar.line.fill.background()

        # Label
        label_box = slide.shapes.add_textbox(
            x, y + height + Inches(0.1), width, Inches(0.4)
        )
        label_para = label_box.text_frame.paragraphs[0]
        label_para.text = f"{label}: {percentage:.0f}%"
        label_para.alignment = PP_ALIGN.CENTER
        label_para.font.size = Pt(12)
        text_color = colors.get('text', '#2C3E50').lstrip('#')
        label_para.font.color.rgb = RGBColor.from_string(text_color)


class SmartLayoutDetector:
    """Detects optimal layout pattern based on content analysis."""

    @staticmethod
    def detect_layout(title: str, content_items: List[str]) -> str:
        """Detect the best layout for this slide's content."""
        if not content_items:
            return "simple"

        title_lower = title.lower()
        content_text = ' '.join(content_items).lower()

        # Quote detection: single item under 30 words
        if len(content_items) == 1 and len(content_items[0].split()) < 30:
            return "quote"

        # Timeline detection
        if any(word in title_lower for word in ['timeline', 'phases', 'roadmap', 'schedule']):
            return "timeline"
        if any(word in content_text for word in ['q1', 'q2', 'q3', 'q4', 'phase']):
            return "timeline"

        # Hierarchy detection
        if any(word in title_lower for word in ['architecture', 'structure', 'hierarchy', 'organization']):
            return "hierarchy"

        # Comparison detection
        if any(word in title_lower for word in ['vs', 'comparison', 'versus']):
            return "comparison"
        if any(word in content_text for word in ['versus', 'compared to', 'current', 'proposed']):
            return "comparison"

        # Metrics/data detection
        has_percentages = any('%' in item for item in content_items)
        if any(word in title_lower for word in ['metrics', 'kpi', 'results', 'performance']):
            return "metrics"
        if has_percentages and len(content_items) >= 2:
            return "metrics"

        # Table detection: items with pipe/tab delimiters
        if any('|' in item or '\t' in item for item in content_items):
            return "table"

        # Grid layout for 4-6 items
        if 4 <= len(content_items) <= 6:
            return "grid"

        # Multi-column if many items
        if len(content_items) > 6:
            return "multi_column"

        return "simple"


class AIEnhancedPPTXRenderer:
    """AI-Enhanced PPTX renderer with Hugging Face design intelligence"""

    def __init__(self, huggingface_api_key: Optional[str] = None):
        self.ai_generator = None
        self.background_renderer = CompositeBackgroundRenderer()
        self.chart_generator = DynamicChartGenerator()
        self.layout_detector = SmartLayoutDetector()
        self.content_analyzer = ContentAnalyzer()
        self.design_rules = SmartDesignRules()
        self.layout_engine = AdvancedLayoutEngine()
        self.logger = logging.getLogger(__name__)

        # Initialize AI generator if API key is available
        if huggingface_api_key:
            self.ai_generator = AIDesignGenerator(huggingface_api_key)
            self.logger.info("AI design generator initialized with Hugging Face")
        else:
            self.logger.warning("No Hugging Face API key - using default themes only")

    def hex_to_rgb(self, hex_color):
        """Convert hex color to RGBColor"""
        hex_color = hex_color.lstrip('#')
        return RGBColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16)
        )

    @staticmethod
    def _geometry_to_inches(value, slide_dimension):
        """Convert geometry value to inches. Values <= 1.0 are treated as fractions."""
        if value <= 1.0:
            return value * slide_dimension
        return value

    async def analyze_content_with_ai(self, json_data: Dict[str, Any], company_info: Dict[str, Any]) -> Optional[Dict[str, Any]]:
        """Use AI to analyze content and generate design recommendations"""
        if not self.ai_generator:
            return None

        try:
            self.logger.info("Analyzing content with Hugging Face AI...")
            design_analysis = await self.ai_generator.analyze_content_for_unique_design(json_data, company_info)
            self.logger.info(f"AI analysis complete: {design_analysis.get('industry')} industry, {design_analysis.get('style')} style")
            return design_analysis
        except Exception as e:
            self.logger.error(f"AI analysis failed: {e}")
            return None

    def get_design_theme(self, ai_analysis: Optional[Dict[str, Any]], fallback_content: str = "", spec_colors: Optional[Dict[str, str]] = None) -> Any:
        """Get design theme based on AI analysis or fallback to industry detection"""
        if ai_analysis:
            industry = ai_analysis.get('industry', 'corporate')
            theme = DesignTemplateLibrary.get_theme_for_industry(industry)
            self.logger.info(f"Using AI-determined theme: {theme.name} for {industry}")
        else:
            style_analysis = {'industry': self._detect_industry_from_content(fallback_content)}
            design_system = get_design_system_for_content(fallback_content, style_analysis)
            theme = design_system['theme']
            self.logger.info(f"Using fallback theme: {theme.name}")

        # Override theme colors with spec's tokens.colors if provided
        if spec_colors:
            self.logger.info(f"Applying spec color overrides: {list(spec_colors.keys())}")
            for key, value in spec_colors.items():
                if isinstance(value, str) and value.startswith('#') and len(value) == 7:
                    theme.colors[key] = value

        return theme

    def _detect_industry_from_content(self, content: str) -> str:
        """Basic industry detection from content text"""
        content_lower = content.lower()
        if any(word in content_lower for word in ['health', 'medical', 'patient', 'hospital']):
            return 'healthcare'
        elif any(word in content_lower for word in ['finance', 'bank', 'investment', 'money']):
            return 'finance'
        elif any(word in content_lower for word in ['tech', 'software', 'digital', 'api']):
            return 'technology'
        elif any(word in content_lower for word in ['security', 'cyber', 'protection']):
            return 'security'
        elif any(word in content_lower for word in ['education', 'learning', 'training']):
            return 'education'
        else:
            return 'corporate'

    def apply_ai_enhanced_background(self, slide, design_theme):
        """Apply AI-enhanced background using olama's background renderers"""
        design_config = {
            'background_design': design_theme.background_design,
            'watermark': design_theme.watermark,
            'colors': design_theme.colors
        }

        try:
            self.background_renderer.render_background(slide, design_config)
            self.logger.debug(f"Applied {design_theme.name} background")
        except Exception as e:
            self.logger.warning(f"Background rendering failed: {e}")
            # Apply simple background as fallback
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = self.hex_to_rgb(design_theme.colors['background'])

    def add_text_with_design_theme(self, slide, content, x, y, width, height, text_type, design_theme,
                                    font_size_override: Optional[int] = None):
        """Add text with intelligent theme-based styling from olama design system"""
        text_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.clear()

        p = text_frame.paragraphs[0]

        # Handle multi-line content with bullets
        lines = content.split('\n')
        for i, line in enumerate(lines):
            if i > 0:
                p = text_frame.add_paragraph()
                p.text = f"  \u2022  {line.strip()}" if line.strip() else ""
            else:
                p.text = line.strip()

            # Apply typography from design theme
            if text_type == 'title':
                typography = design_theme.typography.get('title_slide', {})
            elif text_type == 'subtitle':
                typography = design_theme.typography.get('slide_title', {})
            else:
                typography = design_theme.typography.get('body_text', {})

            # Apply font settings
            p.font.name = typography.get('font_name', 'Calibri')
            p.font.size = Pt(font_size_override or typography.get('font_size', 14))
            p.font.bold = typography.get('bold', False)

            # Apply color from design theme, with contrast checking
            color_key = typography.get('color', 'text')
            if color_key in design_theme.colors:
                text_color = design_theme.colors[color_key]
                bg_color = design_theme.colors.get('background', '#FFFFFF')
                safe_color = self.design_rules.ensure_contrast_ratio(bg_color, text_color)
                p.font.color.rgb = self.hex_to_rgb(safe_color)

            if len(lines) > 1 and i > 0:
                p.space_before = Pt(6)

    def _render_title(self, slide, title_text: str, title_ph, slide_w, slide_h, design_theme):
        """Render the slide title using geometry from placeholder or defaults."""
        if title_ph and 'geometry' in title_ph:
            g = title_ph['geometry']
            x = self._geometry_to_inches(g.get('x', 0.05), slide_w)
            y = self._geometry_to_inches(g.get('y', 0.05), slide_w)
            w = self._geometry_to_inches(g.get('w', 0.9), slide_w)
            h = self._geometry_to_inches(g.get('h', 0.12), slide_h)
        else:
            x, y, w, h = 0.5, 0.5, slide_w - 1.0, 1.0
        self.add_text_with_design_theme(slide, title_text, x, y, w, h, 'title', design_theme)

    def _render_simple_layout(self, slide, items: List[str], slide_w, slide_h, design_theme):
        """Render body items as a simple bulleted list."""
        if not items:
            return
        content = '\n'.join(items)
        x, y = 0.8, 1.8
        w, h = slide_w - 1.6, slide_h - 2.5
        self.add_text_with_design_theme(slide, content, x, y, w, h, 'body', design_theme)

    def _render_metrics_layout(self, slide, items: List[str], slide_w, slide_h, design_theme):
        """Render metrics/data with charts or progress bars when data is detected."""
        colors = design_theme.colors

        # Try to detect chartable data
        chart_data = self.chart_generator.detect_data_pattern(items)

        if chart_data["chart_type"] and len(chart_data["data"]) >= 2:
            # Create chart on left, supporting text on right
            chart_x = Inches(0.5)
            chart_y = Inches(2.0)
            chart_w = Inches(slide_w * 0.55)
            chart_h = Inches(slide_h * 0.55)

            self.chart_generator.create_chart(
                slide, chart_data, chart_x, chart_y, chart_w, chart_h
            )

            # Add remaining text on the right
            non_chart_items = [item for item in items
                               if not any(label in item for label in chart_data["labels"])]
            if non_chart_items:
                text_content = '\n'.join(non_chart_items)
                text_x = slide_w * 0.6
                self.add_text_with_design_theme(
                    slide, text_content, text_x, 2.0,
                    slide_w * 0.35, slide_h * 0.55, 'body', design_theme
                )
            self.logger.info(f"Created {chart_data['chart_type']} chart with {len(chart_data['data'])} data points")

        elif chart_data["has_percentages"] and len(chart_data["data"]) == 1:
            # Single percentage → progress bar
            pct = chart_data["data"][0]
            label = chart_data["labels"][0] if chart_data["labels"] else "Progress"

            self.chart_generator.create_progress_bar(
                slide, Inches(1), Inches(3.0),
                Inches(slide_w - 2), Inches(0.5),
                pct, label, colors
            )

            # Remaining items below
            remaining = [item for item in items if str(int(pct)) + '%' not in item]
            if remaining:
                content = '\n'.join(remaining)
                self.add_text_with_design_theme(
                    slide, content, 1.0, 4.5,
                    slide_w - 2, slide_h - 5.0, 'body', design_theme
                )
            self.logger.info(f"Created progress bar: {pct}%")

        else:
            # No chartable data → metrics grid (2x2 or 3x2)
            self._render_metrics_grid(slide, items, slide_w, slide_h, design_theme)

    def _render_metrics_grid(self, slide, items: List[str], slide_w, slide_h, design_theme):
        """Render items in a grid layout for metrics/KPIs."""
        colors = design_theme.colors
        n = len(items)
        cols = 2 if n <= 4 else 3
        rows = (n + cols - 1) // cols

        grid_w = slide_w - 1.0
        grid_h = slide_h - 3.0
        cell_w = grid_w / cols
        cell_h = grid_h / rows
        start_x, start_y = 0.5, 2.0

        for i, item in enumerate(items):
            if i >= cols * rows:
                break
            row = i // cols
            col = i % cols
            x = start_x + col * cell_w
            y = start_y + row * cell_h

            box = slide.shapes.add_textbox(
                Inches(x), Inches(y),
                Inches(cell_w * 0.9), Inches(cell_h * 0.8)
            )
            frame = box.text_frame
            frame.word_wrap = True
            para = frame.paragraphs[0]
            para.text = str(item).strip()
            para.alignment = PP_ALIGN.CENTER

            typography = design_theme.typography.get('body_text', {})
            para.font.name = typography.get('font_name', 'Calibri')
            para.font.size = Pt(typography.get('font_size', 14) + 2)
            para.font.bold = True
            color_key = typography.get('color', 'text')
            if color_key in colors:
                para.font.color.rgb = self.hex_to_rgb(colors[color_key])

    def _render_timeline_layout(self, slide, items: List[str], slide_w, slide_h, design_theme):
        """Render items as horizontal timeline cards."""
        if not items:
            return
        colors = design_theme.colors
        n = min(len(items), 6)  # Max 6 timeline cards
        card_w = (slide_w - 1.5) / n
        card_h = 2.5
        start_x, start_y = 0.5, 2.5

        # Draw a connecting line
        line = slide.shapes.add_connector(
            1,  # STRAIGHT
            Inches(start_x), Inches(start_y + card_h / 2),
            Inches(start_x + n * card_w), Inches(start_y + card_h / 2)
        )
        primary = colors.get('primary', '#2E75B6')
        line.line.color.rgb = self.hex_to_rgb(primary)
        line.line.width = Pt(2)

        for i in range(n):
            x = start_x + i * card_w + 0.1
            # Card background
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(start_y),
                Inches(card_w - 0.2), Inches(card_h)
            )
            light = colors.get('light', '#F8F9FA')
            card.fill.solid()
            card.fill.fore_color.rgb = self.hex_to_rgb(light)
            card.line.color.rgb = self.hex_to_rgb(primary)
            card.line.width = Pt(1)

            # Card text
            text_box = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(start_y + 0.2),
                Inches(card_w - 0.4), Inches(card_h - 0.4)
            )
            frame = text_box.text_frame
            frame.word_wrap = True
            para = frame.paragraphs[0]
            para.text = items[i].strip()
            para.alignment = PP_ALIGN.CENTER

            typography = design_theme.typography.get('body_text', {})
            para.font.name = typography.get('font_name', 'Calibri')
            para.font.size = Pt(min(typography.get('font_size', 14), 13))
            color_key = typography.get('color', 'text')
            if color_key in colors:
                para.font.color.rgb = self.hex_to_rgb(colors[color_key])

        # Remaining items below if any
        if len(items) > n:
            extra = '\n'.join(items[n:])
            self.add_text_with_design_theme(
                slide, extra, 0.5, start_y + card_h + 0.5,
                slide_w - 1.0, slide_h - start_y - card_h - 1.0, 'body', design_theme
            )

    def _render_comparison_layout(self, slide, items: List[str], slide_w, slide_h, design_theme):
        """Render items in two columns for comparison."""
        if not items:
            return
        colors = design_theme.colors
        mid = len(items) // 2
        left_items = items[:mid] if mid > 0 else items[:1]
        right_items = items[mid:] if mid > 0 else items[1:]

        col_w = (slide_w - 1.5) / 2
        start_y = 2.0

        # Left column header bar
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(start_y - 0.4),
            Inches(col_w), Inches(0.35)
        )
        primary = colors.get('primary', '#2E75B6')
        header.fill.solid()
        header.fill.fore_color.rgb = self.hex_to_rgb(primary)
        header.line.fill.background()

        # Right column header bar
        header2 = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5 + col_w + 0.5), Inches(start_y - 0.4),
            Inches(col_w), Inches(0.35)
        )
        accent = colors.get('accent', '#3498DB')
        header2.fill.solid()
        header2.fill.fore_color.rgb = self.hex_to_rgb(accent)
        header2.line.fill.background()

        # Left column content
        left_content = '\n'.join(left_items)
        self.add_text_with_design_theme(
            slide, left_content, 0.5, start_y,
            col_w, slide_h - start_y - 1.0, 'body', design_theme
        )

        # Right column content
        right_content = '\n'.join(right_items)
        self.add_text_with_design_theme(
            slide, right_content, 0.5 + col_w + 0.5, start_y,
            col_w, slide_h - start_y - 1.0, 'body', design_theme
        )

    def _render_multi_column_layout(self, slide, items: List[str], slide_w, slide_h, design_theme):
        """Render items split across 2 columns."""
        if not items:
            return
        mid = len(items) // 2
        left_items = items[:mid]
        right_items = items[mid:]

        col_w = (slide_w - 1.5) / 2
        start_y = 2.0

        left_content = '\n'.join(left_items)
        self.add_text_with_design_theme(
            slide, left_content, 0.5, start_y,
            col_w, slide_h - start_y - 1.0, 'body', design_theme
        )

        right_content = '\n'.join(right_items)
        self.add_text_with_design_theme(
            slide, right_content, 0.5 + col_w + 0.5, start_y,
            col_w, slide_h - start_y - 1.0, 'body', design_theme
        )

    def _render_quote_layout(self, slide, items: List[str], slide_w, slide_h, design_theme):
        """Render a single item as a large centered quote."""
        if not items:
            return
        colors = design_theme.colors
        quote_text = items[0].strip()

        # Large centered quote text
        margin_x = slide_w * 0.15
        text_w = slide_w - 2 * margin_x
        text_h = slide_h * 0.4
        text_y = (slide_h - text_h) / 2  # Vertically centered

        text_box = slide.shapes.add_textbox(
            Inches(margin_x), Inches(text_y),
            Inches(text_w), Inches(text_h)
        )
        frame = text_box.text_frame
        frame.word_wrap = True
        para = frame.paragraphs[0]
        para.text = f'\u201C{quote_text}\u201D'
        para.alignment = PP_ALIGN.CENTER

        # Use larger body font for quotes
        typography = design_theme.typography.get('body_text', {})
        para.font.name = typography.get('font_name', 'Georgia')
        para.font.size = Pt(28)
        para.font.italic = True
        primary = colors.get('primary', '#2E75B6')
        para.font.color.rgb = self.hex_to_rgb(primary)

        # Decorative accent line below quote
        line_w = slide_w * 0.3
        line_x = (slide_w - line_w) / 2
        line_y = text_y + text_h + 0.3
        accent_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(line_x), Inches(line_y),
            Inches(line_w), Inches(0.06)
        )
        accent = colors.get('accent', '#3498DB')
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = self.hex_to_rgb(accent)
        accent_line.line.fill.background()

    def _render_grid_layout(self, slide, items: List[str], slide_w, slide_h, design_theme):
        """Render 4-6 items in a 2x2 or 2x3 grid with card backgrounds."""
        if not items:
            return
        colors = design_theme.colors
        n = min(len(items), 6)
        cols = 2 if n <= 4 else 3
        rows = (n + cols - 1) // cols

        grid_w = slide_w - 1.0
        grid_h = slide_h - 3.0
        cell_w = grid_w / cols
        cell_h = grid_h / rows
        start_x, start_y = 0.5, 2.0
        padding = 0.15

        for i in range(n):
            row = i // cols
            col = i % cols
            x = start_x + col * cell_w + padding
            y = start_y + row * cell_h + padding
            w = cell_w - 2 * padding
            h = cell_h - 2 * padding

            # Card background
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y), Inches(w), Inches(h)
            )
            light = colors.get('light', '#F8F9FA')
            card.fill.solid()
            card.fill.fore_color.rgb = self.hex_to_rgb(light)
            primary = colors.get('primary', '#2E75B6')
            card.line.color.rgb = self.hex_to_rgb(primary)
            card.line.width = Pt(1)

            # Card text
            text_box = slide.shapes.add_textbox(
                Inches(x + 0.15), Inches(y + 0.15),
                Inches(w - 0.3), Inches(h - 0.3)
            )
            frame = text_box.text_frame
            frame.word_wrap = True
            para = frame.paragraphs[0]
            para.text = items[i].strip()
            para.alignment = PP_ALIGN.CENTER

            typography = design_theme.typography.get('body_text', {})
            para.font.name = typography.get('font_name', 'Calibri')
            para.font.size = Pt(typography.get('font_size', 14))
            para.font.bold = True
            color_key = typography.get('color', 'text')
            if color_key in colors:
                para.font.color.rgb = self.hex_to_rgb(colors[color_key])

    def _create_table(self, slide, headers: List[str], rows: List[List[str]],
                       x: float, y: float, width: float, height: float,
                       design_theme) -> None:
        """Create a styled table on the slide (ported from olama).

        Args:
            slide: pptx slide object
            headers: Column header strings
            rows: List of row data (each row is a list of cell strings)
            x, y, width, height: Position and dimensions in inches
            design_theme: DesignTheme for styling
        """
        if not headers or not rows:
            return

        num_rows = len(rows) + 1  # +1 for header
        num_cols = len(headers)
        colors = design_theme.colors

        table_shape = slide.shapes.add_table(
            num_rows, num_cols,
            Inches(x), Inches(y), Inches(width), Inches(height)
        )
        table = table_shape.table

        # Style header row
        for col_idx, header_text in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = str(header_text)
            # Header background
            cell.fill.solid()
            primary = colors.get('primary', '#2E75B6')
            cell.fill.fore_color.rgb = self.hex_to_rgb(primary)
            # Header text styling
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(12)
            para.font.bold = True
            para.font.color.rgb = self.hex_to_rgb('#FFFFFF')
            para.alignment = PP_ALIGN.CENTER

        # Style data rows
        for row_idx, row_data in enumerate(rows):
            for col_idx, cell_text in enumerate(row_data[:num_cols]):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(cell_text)
                # Alternating row colors
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    light = colors.get('light', '#F8F9FA')
                    cell.fill.fore_color.rgb = self.hex_to_rgb(light)
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(11)
                text_color = colors.get('text', '#2C3E50')
                para.font.color.rgb = self.hex_to_rgb(text_color)

    def _add_image(self, slide, image_path: str, x: float, y: float,
                    width: float, height: float) -> bool:
        """Add an image to the slide (ported from olama).

        Args:
            slide: pptx slide object
            image_path: Path to image file
            x, y, width, height: Position and dimensions in inches

        Returns:
            True if image was added, False if file not found
        """
        try:
            if not os.path.exists(image_path):
                self.logger.warning(f"Image not found: {image_path}")
                return False
            slide.shapes.add_picture(
                image_path, Inches(x), Inches(y), Inches(width), Inches(height)
            )
            return True
        except Exception as e:
            self.logger.warning(f"Failed to add image {image_path}: {e}")
            return False

    def _add_slide_number(self, slide, slide_number: int, total_slides: int,
                           slide_w: float, slide_h: float, design_theme) -> None:
        """Add slide number to bottom-right corner."""
        colors = design_theme.colors
        text_box = slide.shapes.add_textbox(
            Inches(slide_w - 1.5), Inches(slide_h - 0.5),
            Inches(1.0), Inches(0.3)
        )
        para = text_box.text_frame.paragraphs[0]
        para.text = f"{slide_number} / {total_slides}"
        para.alignment = PP_ALIGN.RIGHT
        para.font.size = Pt(9)
        secondary = colors.get('secondary', '#666666')
        para.font.color.rgb = self.hex_to_rgb(secondary)

    def _render_table_layout(self, slide, items: List[str], slide_w, slide_h, design_theme):
        """Render items as a table if they contain delimiter-separated data."""
        # Detect table structure from content
        delimiter = None
        for d in ['|', '\t', '  ']:
            if any(d in item for item in items):
                delimiter = d
                break

        if delimiter:
            # Parse table data
            parsed_rows = [item.split(delimiter) for item in items]
            parsed_rows = [[cell.strip() for cell in row if cell.strip()] for row in parsed_rows]
            if parsed_rows and len(parsed_rows) > 1:
                headers = parsed_rows[0]
                rows = parsed_rows[1:]
                self._create_table(
                    slide, headers, rows,
                    0.5, 2.0, slide_w - 1.0, min(slide_h - 3.0, len(rows) * 0.5 + 0.8),
                    design_theme
                )
                return

        # Fallback: render as grid
        self._render_grid_layout(slide, items, slide_w, slide_h, design_theme)

    def _render_hierarchy_layout(self, slide, items: List[str], slide_w, slide_h, design_theme):
        """Render items as a top-down hierarchy (tree structure)."""
        if not items:
            return
        colors = design_theme.colors
        n = len(items)

        # Top node (first item) centered
        top_w, top_h = 3.0, 0.8
        top_x = (slide_w - top_w) / 2
        top_y = 2.0

        top_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(top_x), Inches(top_y), Inches(top_w), Inches(top_h)
        )
        primary = colors.get('primary', '#2E75B6')
        top_box.fill.solid()
        top_box.fill.fore_color.rgb = self.hex_to_rgb(primary)
        top_box.line.fill.background()

        top_text = slide.shapes.add_textbox(
            Inches(top_x + 0.1), Inches(top_y + 0.1),
            Inches(top_w - 0.2), Inches(top_h - 0.2)
        )
        frame = top_text.text_frame
        frame.word_wrap = True
        para = frame.paragraphs[0]
        para.text = items[0].strip()
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(16)
        para.font.bold = True
        para.font.color.rgb = self.hex_to_rgb('#FFFFFF')

        # Child nodes below (remaining items)
        children = items[1:]
        if not children:
            return

        child_count = min(len(children), 5)
        child_w = min(2.5, (slide_w - 1.0) / child_count - 0.2)
        child_h = 0.7
        child_y = top_y + top_h + 1.2
        total_children_w = child_count * child_w + (child_count - 1) * 0.2
        child_start_x = (slide_w - total_children_w) / 2

        # Draw connecting lines from top to children
        top_center_x = top_x + top_w / 2
        top_bottom_y = top_y + top_h

        for i in range(child_count):
            cx = child_start_x + i * (child_w + 0.2)
            child_center_x = cx + child_w / 2

            # Connector line
            connector = slide.shapes.add_connector(
                1, Inches(top_center_x), Inches(top_bottom_y),
                Inches(child_center_x), Inches(child_y)
            )
            connector.line.color.rgb = self.hex_to_rgb(primary)
            connector.line.width = Pt(1.5)

            # Child box
            child_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(cx), Inches(child_y), Inches(child_w), Inches(child_h)
            )
            light = colors.get('light', '#F8F9FA')
            child_box.fill.solid()
            child_box.fill.fore_color.rgb = self.hex_to_rgb(light)
            child_box.line.color.rgb = self.hex_to_rgb(primary)
            child_box.line.width = Pt(1)

            # Child text
            child_text = slide.shapes.add_textbox(
                Inches(cx + 0.05), Inches(child_y + 0.05),
                Inches(child_w - 0.1), Inches(child_h - 0.1)
            )
            cframe = child_text.text_frame
            cframe.word_wrap = True
            cpara = cframe.paragraphs[0]
            cpara.text = children[i].strip()
            cpara.alignment = PP_ALIGN.CENTER
            cpara.font.size = Pt(12)
            color_key = design_theme.typography.get('body_text', {}).get('color', 'text')
            if color_key in colors:
                cpara.font.color.rgb = self.hex_to_rgb(colors[color_key])

        # Remaining items below hierarchy as simple text
        if len(children) > child_count:
            extra = '\n'.join(children[child_count:])
            self.add_text_with_design_theme(
                slide, extra, 0.5, child_y + child_h + 0.5,
                slide_w - 1.0, slide_h - child_y - child_h - 1.0, 'body', design_theme
            )

    async def render_with_ai_design(self, spec_data, output_path, company_info: Optional[Dict[str, Any]] = None):
        """Main rendering method with AI design analysis"""
        self.logger.info("Starting AI-enhanced PPTX rendering...")

        # Extract spec colors from tokens.colors
        spec_colors = None
        tokens = spec_data.get('tokens', {})
        if isinstance(tokens, dict):
            spec_colors = tokens.get('colors')

        # Extract content for analysis
        all_content = self._extract_all_content(spec_data)

        # Perform AI analysis if available
        ai_analysis = None
        if self.ai_generator and company_info:
            json_data = {'slides': self._convert_spec_to_slides(spec_data)}
            ai_analysis = await self.analyze_content_with_ai(json_data, company_info)

        # Get design theme based on AI analysis or fallback, with spec color overrides
        design_theme = self.get_design_theme(ai_analysis, all_content, spec_colors)

        # Create presentation
        prs = Presentation()

        # Remove default slide if any
        if len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]

        # Get slide dimensions for geometry conversion
        slide_w = prs.slide_width / Emu(914400)  # Convert EMU to inches
        slide_h = prs.slide_height / Emu(914400)

        # Process layouts with AI-enhanced design
        layouts = spec_data.get('layouts', [])
        for layout in layouts:
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

            # Apply AI-enhanced background
            self.apply_ai_enhanced_background(slide, design_theme)

            # Extract title and body content from placeholders
            placeholders = layout.get('placeholders', [])
            slide_title = ""
            body_items = []
            title_ph = None

            for ph in placeholders:
                ph_id = ph.get('id', '')
                ph_type = ph.get('type', 'body')
                content = ph.get('content', '')

                if 'subtitle' in ph_id.lower() or 'subheading' in ph_id.lower() or ph_type == 'subtitle':
                    body_items.insert(0, content)  # Subtitle goes first in body
                elif 'title' in ph_id.lower() or 'heading' in ph_id.lower() or ph_type == 'title':
                    slide_title = content
                    title_ph = ph
                else:
                    # Split multi-line body content into separate items
                    if '\n' in content:
                        body_items.extend([line.strip() for line in content.split('\n') if line.strip()])
                    else:
                        body_items.append(content)

            # Content analysis for dynamic font sizing
            analysis = self.content_analyzer.analyze_content(slide_title, body_items)
            font_sizes = self.design_rules.calculate_optimal_font_sizes(
                analysis["visual_weight"], analysis["word_count"], analysis["content_type"]
            )

            # Use explicit layout_hint from spec if it's a known type, otherwise auto-detect
            KNOWN_LAYOUTS = {"title", "quote", "timeline", "hierarchy", "comparison",
                             "metrics", "table", "grid", "multi_column", "simple"}
            layout_name = layout.get('name', '').lower()
            if layout_name in KNOWN_LAYOUTS:
                layout_type = layout_name
            else:
                layout_type = self.layout_detector.detect_layout(slide_title, body_items)
            self.logger.info(
                f"Slide '{slide_title}': layout={layout_type}, "
                f"type={analysis['content_type']}, weight={analysis['visual_weight']:.2f}, "
                f"font title={font_sizes['title']}pt body={font_sizes['body']}pt"
            )

            # Always render the title (with dynamic font size)
            if slide_title:
                self._render_title(slide, slide_title, title_ph, slide_w, slide_h, design_theme)

            # Render body content using smart layout
            if layout_type == "quote":
                self._render_quote_layout(slide, body_items, slide_w, slide_h, design_theme)
            elif layout_type == "hierarchy":
                self._render_hierarchy_layout(slide, body_items, slide_w, slide_h, design_theme)
            elif layout_type == "grid":
                self._render_grid_layout(slide, body_items, slide_w, slide_h, design_theme)
            elif layout_type == "table":
                self._render_table_layout(slide, body_items, slide_w, slide_h, design_theme)
            elif layout_type == "metrics":
                self._render_metrics_layout(slide, body_items, slide_w, slide_h, design_theme)
            elif layout_type == "timeline":
                self._render_timeline_layout(slide, body_items, slide_w, slide_h, design_theme)
            elif layout_type == "comparison":
                self._render_comparison_layout(slide, body_items, slide_w, slide_h, design_theme)
            elif layout_type == "multi_column":
                self._render_multi_column_layout(slide, body_items, slide_w, slide_h, design_theme)
            else:
                self._render_simple_layout(slide, body_items, slide_w, slide_h, design_theme)

        # Add slide numbers
        total_slides = len(prs.slides)
        for i, slide in enumerate(prs.slides):
            self._add_slide_number(slide, i + 1, total_slides, slide_w, slide_h, design_theme)

        # Save presentation
        prs.save(output_path)

        # Log design decisions
        if ai_analysis:
            self.logger.info(f"AI-enhanced presentation generated: {design_theme.name} theme")
            self.logger.info(f"AI insights: {ai_analysis.get('reasoning', 'N/A')}")
        else:
            self.logger.info(f"Presentation generated with {design_theme.name} theme (no AI analysis)")

    def _extract_all_content(self, spec_data) -> str:
        """Extract all text content from spec for analysis"""
        content_parts = []
        layouts = spec_data.get('layouts', [])
        for layout in layouts:
            placeholders = layout.get('placeholders', [])
            for ph in placeholders:
                content = ph.get('content', '')
                if content:
                    content_parts.append(str(content))
        return ' '.join(content_parts)

    def _convert_spec_to_slides(self, spec_data) -> list:
        """Convert spec format to slides format for AI analysis"""
        slides = []
        layouts = spec_data.get('layouts', [])
        for layout in layouts:
            slide_content = []
            slide_title = ""
            placeholders = layout.get('placeholders', [])
            for ph in placeholders:
                content = ph.get('content', '')
                ph_type = ph.get('type', 'body')
                if ph_type == 'title':
                    slide_title = content
                else:
                    slide_content.append(content)
            slides.append({'title': slide_title, 'content': slide_content})
        return slides

    def render_pptx_sync(self, spec_data, output_path, company_info: Optional[Dict[str, Any]] = None):
        """Synchronous wrapper for backward compatibility"""
        return asyncio.run(self.render_with_ai_design(spec_data, output_path, company_info))


async def main():
    parser = argparse.ArgumentParser(description='AI-Enhanced PPTX Renderer with Hugging Face')
    parser.add_argument('spec_file', help='JSON spec file')
    parser.add_argument('output_file', help='Output PPTX file')
    parser.add_argument('--company-info', help='Company info JSON file (optional)')
    parser.add_argument('--hf-api-key', help='Hugging Face API key (or set HUGGING_FACE_API_KEY env var)')

    args = parser.parse_args()

    try:
        # Load spec
        with open(args.spec_file, 'r') as f:
            spec_data = json.load(f)

        # Load company info if provided
        company_info = None
        if args.company_info:
            with open(args.company_info, 'r') as f:
                company_info = json.load(f)

        # Get API key
        api_key = args.hf_api_key or os.getenv('HUGGING_FACE_API_KEY')

        # Render with AI enhancement
        renderer = AIEnhancedPPTXRenderer(api_key)
        await renderer.render_with_ai_design(spec_data, args.output_file, company_info)

        print(f"Generated: {args.output_file}")

    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)

def main_sync():
    """Synchronous entry point for backward compatibility"""
    asyncio.run(main())


if __name__ == '__main__':
    main_sync()
