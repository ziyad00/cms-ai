"""AI-powered design generator for PowerPoint presentations using Hugging Face AI."""
import json
import logging
import asyncio
import os
from typing import Dict, List, Any, Optional, Tuple
from datetime import datetime
import httpx

logger = logging.getLogger(__name__)

# Import mock generator for development/testing
from mock_ai_generator import get_mock_generator


class AIDesignGenerator:
    """Generate professional PowerPoint designs using Hugging Face AI."""

    def __init__(self, api_key: str = None, model: str = "mistralai/Mixtral-8x7B-Instruct-v0.1"):
        """
        Initialize the AI design generator.

        Args:
            api_key: Hugging Face API key (optional, uses mock if not provided)
            model: AI model to use (default: Mixtral-8x7B)
        """
        self.api_key = api_key
        self.model = model
        self.base_url = "https://router.huggingface.co/v1/chat/completions"
        self.client = httpx.AsyncClient()

        # Use mock if configured or no API key
        self.use_mock = os.getenv('USE_MOCK_AI', 'false').lower() == 'true' or not api_key
        if self.use_mock:
            self.mock_generator = get_mock_generator()
            logger.info("Using mock AI generator (no API costs)")

    async def analyze_content_for_unique_design(self, json_data: Dict[str, Any], company_info: Dict[str, Any]) -> Dict[str, str]:
        """Analyze content deeply to generate truly unique design recommendations."""
        # Use mock if configured
        if self.use_mock:
            return await self.mock_generator.analyze_content_for_unique_design(json_data, company_info)

        # Extract all text content for analysis
        all_content = []
        for slide in json_data.get('slides', []):
            all_content.append(slide.get('title', ''))
            all_content.extend(slide.get('content', []))

        content_text = ' '.join([str(item) for item in all_content])

        # Count key themes and concepts
        tech_keywords = ['api', 'database', 'architecture', 'backend', 'frontend', 'cloud', 'digital']
        business_keywords = ['strategy', 'governance', 'stakeholder', 'management', 'roi', 'value']
        security_keywords = ['security', 'encryption', 'compliance', 'risk', 'authentication']
        innovation_keywords = ['ai', 'machine learning', 'innovation', 'automation', 'future']

        keyword_counts = {
            'tech': sum(1 for word in tech_keywords if word in content_text.lower()),
            'business': sum(1 for word in business_keywords if word in content_text.lower()),
            'security': sum(1 for word in security_keywords if word in content_text.lower()),
            'innovation': sum(1 for word in innovation_keywords if word in content_text.lower())
        }

        # Determine dominant theme
        dominant_theme = max(keyword_counts, key=keyword_counts.get)
        theme_strength = keyword_counts[dominant_theme]

        prompt = f"""
        CONTENT ANALYSIS:
        - Total slides: {len(json_data.get('slides', []))}
        - Dominant theme: {dominant_theme} (strength: {theme_strength})
        - Content complexity: {"high" if len(content_text.split()) > 1000 else "medium" if len(content_text.split()) > 500 else "low"}
        - Key concepts detected: {', '.join(set([word for word in content_text.split() if len(word) > 6][:10]))}

        COMPANY CONTEXT:
        {json.dumps(company_info, indent=2)}

        SAMPLE CONTENT (first 3 slides):
        {content_text[:800]}...

        Based on this SPECIFIC content and context, create a completely unique visual design identity.
        Don't use generic styles - create something that perfectly captures THIS presentation's essence.

        Consider:
        - What emotions should this content evoke?
        - What visual metaphors represent these concepts?
        - How formal vs innovative should the design feel?
        - What industries and audiences are involved?

        Return UNIQUE design recommendations as JSON:
        {{
            "industry": "specific industry/sector",
            "formality": "exact formality level",
            "style": "unique style description",
            "color_preference": "specific color approach with reasoning",
            "audience": "precise target audience",
            "visual_metaphor": "what visual concepts represent this content",
            "emotional_tone": "what feeling should the design convey",
            "reasoning": "detailed explanation of design choices"
        }}
        """

        try:
            headers = {
                "Authorization": f"Bearer {self.api_key}",
                "Content-Type": "application/json"
            }

            payload = {
                "model": self.model,
                "messages": [
                    {"role": "system", "content": "You are a expert design strategist who creates unique visual identities based on content analysis. Return detailed, specific recommendations, not generic categories."},
                    {"role": "user", "content": prompt}
                ],
                "temperature": 0.7,  # Higher creativity for unique designs
                "max_tokens": 600
            }

            response = await self.client.post(self.base_url, headers=headers, json=payload)
            response.raise_for_status()

            result = response.json()
            message = result["choices"][0]["message"]
            content = message.get("content") or message.get("reasoning_content") or ""
            content = content.strip() if content else ""

            logger.info(f"Raw AI response: {content}")

            # Extract JSON from response
            if '```json' in content:
                content = content.split('```json')[1].split('```')[0].strip()
            elif '{' in content:
                start = content.find('{')
                end = content.rfind('}') + 1
                content = content[start:end]

            return json.loads(content)

        except Exception as e:
            logger.error(f"Error analyzing content style: {e}")
            # Return default professional style
            return {
                "industry": "government",
                "formality": "formal",
                "style": "corporate",
                "color_preference": "professional blue",
                "audience": "executives",
                "visual_metaphor": "structure and reliability",
                "emotional_tone": "trustworthy and professional",
                "reasoning": "fallback to safe professional style"
            }

    async def analyze_content_style(self, json_data: Dict[str, Any], company_info: Dict[str, Any]) -> Dict[str, str]:
        # Use mock if configured
        if self.use_mock:
            return await self.mock_generator.analyze_content_style(json_data, company_info)
        """
        Analyze content to determine appropriate design style.

        Args:
            json_data: Presentation content
            company_info: Company information

        Returns:
            Dictionary with style recommendations
        """
        # Extract key information for analysis
        slides_content = []
        for slide in json_data.get('slides', []):
            title = slide.get('title', '')
            content = ' '.join(slide.get('content', []))
            slides_content.append(f"Title: {title}\nContent: {content}")

        content_text = '\n\n'.join(slides_content[:5])  # First 5 slides for analysis

        prompt = f"""
        Analyze this presentation content and company information to recommend design style:

        COMPANY INFO:
        {json.dumps(company_info, indent=2)}

        PRESENTATION CONTENT (first 5 slides):
        {content_text}

        Based on this content, recommend:
        1. Industry type (government, tech, healthcare, finance, education, etc.)
        2. Formality level (formal, business, casual)
        3. Design style (corporate, modern, minimal, creative)
        4. Color scheme preference (professional blue, tech gradient, healthcare green, etc.)
        5. Target audience (executives, technical teams, public sector, etc.)

        Return as JSON:
        {{
            "industry": "string",
            "formality": "string",
            "style": "string",
            "color_preference": "string",
            "audience": "string",
            "reasoning": "brief explanation"
        }}
        """

        try:
            headers = {
                "Authorization": f"Bearer {self.api_key}",
                "Content-Type": "application/json"
            }

            payload = {
                "model": self.model,
                "messages": [
                    {"role": "system", "content": "You are a professional presentation design expert. Analyze content and recommend appropriate visual design styles."},
                    {"role": "user", "content": prompt}
                ],
                "temperature": 0.3,
                "max_tokens": 500
            }

            response = await self.client.post(self.base_url, headers=headers, json=payload)
            response.raise_for_status()

            result = response.json()
            message = result["choices"][0]["message"]
            content = message.get("content") or message.get("reasoning_content") or ""
            content = content.strip() if content else ""

            logger.info(f"Raw AI response: {content}")

            # Extract JSON from response
            if '```json' in content:
                content = content.split('```json')[1].split('```')[0].strip()
            elif '{' in content:
                start = content.find('{')
                end = content.rfind('}') + 1
                content = content[start:end]

            return json.loads(content)

        except Exception as e:
            logger.error(f"Error analyzing content style: {e}")
            # Return default professional style
            return {
                "industry": "government",
                "formality": "formal",
                "style": "corporate",
                "color_preference": "professional blue",
                "audience": "executives",
                "reasoning": "Default professional style due to analysis error"
            }

    async def generate_color_scheme(self, style_analysis: Dict[str, str]) -> Dict[str, str]:
        # Use mock if configured
        if self.use_mock:
            return await self.mock_generator.generate_color_scheme(style_analysis)
        """
        Generate a professional color scheme based on style analysis.

        Args:
            style_analysis: Results from analyze_content_style

        Returns:
            Dictionary with hex color codes
        """
        prompt = """
        You are a professional color designer. Generate a unique color scheme in valid JSON format.

        Return ONLY valid JSON with these exact keys and hex color values:
        {"primary": "#XXXXXX", "secondary": "#XXXXXX", "background": "#XXXXXX", "text": "#XXXXXX", "accent": "#XXXXXX", "light": "#XXXXXX"}

        Requirements:
        - All hex codes must be 6 characters (e.g., #1A2B3C)
        - Colors must be unique and professional
        - Ensure high contrast between text and background
        - No explanatory text, only JSON
        """

        # Retry mechanism for robust JSON generation
        for attempt in range(3):
            try:
                headers = {
                    "Authorization": f"Bearer {self.api_key}",
                    "Content-Type": "application/json"
                }

                payload = {
                    "model": self.model,
                    "messages": [
                        {"role": "system", "content": "You are a professional color expert. Return ONLY valid JSON, no extra text."},
                        {"role": "user", "content": prompt}
                    ],
                    "temperature": 0.1,
                    "max_tokens": 100
                }

                response = await self.client.post(self.base_url, headers=headers, json=payload)
                response.raise_for_status()

                result = response.json()
                message = result["choices"][0]["message"]
                content = message.get("content") or message.get("reasoning_content") or ""
                content = content.strip() if content else ""

                logger.info(f"AI response attempt {attempt + 1}: {content}")

                # More robust JSON extraction
                import re
                json_content = None

                # Method 1: JSON code blocks
                if '```json' in content:
                    json_content = content.split('```json')[1].split('```')[0].strip()

                # Method 2: Find complete JSON object
                elif '{' in content and '}' in content:
                    # Find the outermost JSON object
                    start = content.find('{')
                    brace_count = 0
                    end = start

                    for i in range(start, len(content)):
                        if content[i] == '{':
                            brace_count += 1
                        elif content[i] == '}':
                            brace_count -= 1
                            if brace_count == 0:
                                end = i + 1
                                break

                    if end > start:
                        json_content = content[start:end]

                # Method 3: Regex for JSON-like pattern
                if not json_content:
                    json_matches = re.findall(r'\{[^{}]*(?:\{[^{}]*\}[^{}]*)*\}', content)
                    if json_matches:
                        json_content = json_matches[-1]  # Take the last/most complete one

                if json_content:
                    try:
                        # Clean JSON content
                        json_content = json_content.strip()
                        # Remove any trailing text after the closing brace
                        if '}' in json_content:
                            json_content = json_content[:json_content.rfind('}') + 1]

                        logger.info(f"Extracted JSON: {json_content}")

                        color_data = json.loads(json_content)

                        # Validate required keys and values
                        required_keys = ["primary", "secondary", "background", "text", "accent", "light"]
                        if all(key in color_data for key in required_keys):
                            # Validate hex color format
                            valid_colors = True
                            for key, color in color_data.items():
                                if not isinstance(color, str) or not re.match(r'^#[0-9A-Fa-f]{6}$', color):
                                    logger.warning(f"Invalid color format for {key}: {color}")
                                    valid_colors = False
                                    break

                            if valid_colors:
                                logger.info(f"Successfully generated color scheme on attempt {attempt + 1}")
                                return color_data
                            else:
                                logger.warning(f"Invalid color format in attempt {attempt + 1}")

                    except json.JSONDecodeError as e:
                        logger.warning(f"JSON decode error on attempt {attempt + 1}: {e}")
                        continue

            except Exception as e:
                logger.warning(f"Color generation attempt {attempt + 1} failed: {e}")
                if attempt == 2:  # Last attempt
                    break
                continue

        # If all attempts failed, use fallback
        logger.error("All color generation attempts failed")
        logger.info("Using fallback color scheme generation")

        # Return default professional blue scheme based on industry
        industry = style_analysis.get('industry', 'government')
        if 'tech' in industry.lower():
            return {
                "primary": "#1E88E5",
                "secondary": "#42A5F5",
                "background": "#FFFFFF",
                "text": "#263238",
                "accent": "#00ACC1",
                "light": "#E3F2FD"
            }
        else:
            return {
                "primary": "#2E75B6",
                "secondary": "#8C8C8C",
                "background": "#FFFFFF",
                "text": "#1F1F1F",
                "accent": "#00B050",
                "light": "#F8F9FA"
            }

    async def generate_typography_system(self, style_analysis: Dict[str, str]) -> Dict[str, Dict[str, Any]]:
        """Generate a typography system using AI based on style analysis (ported from olama).

        Args:
            style_analysis: Results from analyze_content_style

        Returns:
            Typography dict with title_slide, slide_title, body_text, caption
        """
        if self.use_mock:
            return self.mock_generator._get_typography_for_industry(
                style_analysis.get('industry', 'corporate')
            )

        industry = style_analysis.get('industry', 'corporate')
        formality = style_analysis.get('formality', 'professional')
        style = style_analysis.get('style', 'corporate')

        prompt = f"""
        Design a typography system for a {formality} {industry} presentation with {style} style.

        Return ONLY valid JSON with these exact keys:
        {{
            "title_slide": {{"font_name": "FontName", "font_size": 36, "bold": true, "color": "primary"}},
            "slide_title": {{"font_name": "FontName", "font_size": 24, "bold": true, "color": "primary"}},
            "body_text": {{"font_name": "FontName", "font_size": 14, "bold": false, "color": "text"}},
            "caption": {{"font_name": "FontName", "font_size": 11, "bold": false, "color": "secondary"}}
        }}

        Rules:
        - Use professional, widely-available fonts (Calibri, Arial, Segoe UI, Georgia, Times New Roman, Helvetica, Verdana, Garamond, Open Sans, Montserrat, Trebuchet MS)
        - Title fonts should be larger and bolder
        - Body text must be highly readable (14-16pt)
        - Caption should be smaller than body (10-12pt)
        - Font sizes must be integers
        """

        try:
            headers = {
                "Authorization": f"Bearer {self.api_key}",
                "Content-Type": "application/json"
            }
            payload = {
                "model": self.model,
                "messages": [
                    {"role": "system", "content": "You are a typography expert. Return ONLY valid JSON."},
                    {"role": "user", "content": prompt}
                ],
                "temperature": 0.3,
                "max_tokens": 300
            }

            response = await self.client.post(self.base_url, headers=headers, json=payload)
            response.raise_for_status()

            result = response.json()
            content = result["choices"][0]["message"].get("content", "").strip()

            # Extract JSON
            if '```json' in content:
                content = content.split('```json')[1].split('```')[0].strip()
            elif '{' in content:
                start = content.find('{')
                end = content.rfind('}') + 1
                content = content[start:end]

            typography = json.loads(content)

            # Validate structure
            required = ['title_slide', 'body_text']
            if all(k in typography for k in required):
                return typography

        except Exception as e:
            logger.warning(f"Typography generation failed: {e}")

        # Fallback typography
        return self._get_fallback_typography(industry)

    async def generate_slide_layout_code(self, slide_content: Dict[str, Any],
                                          style_analysis: Dict[str, str]) -> Optional[str]:
        """Generate Python code for custom slide layout using AI (ported from olama).

        Args:
            slide_content: Slide title and content items
            style_analysis: Style analysis results

        Returns:
            Python code string or None on failure
        """
        if self.use_mock:
            return self._generate_fallback_slide_code(slide_content)

        title = slide_content.get('title', '')
        content = slide_content.get('content', [])

        prompt = f"""
        Generate python-pptx code to lay out a slide with:
        - Title: "{title}"
        - Content items: {content[:5]}
        - Style: {style_analysis.get('style', 'professional')}
        - Industry: {style_analysis.get('industry', 'corporate')}

        The code should define a function:
        def layout_slide(slide, colors, typography):
            # Add shapes, text boxes, etc.
            pass

        Use python-pptx imports: Inches, Pt, MSO_SHAPE, RGBColor, PP_ALIGN
        Return ONLY the function code, no imports or explanations.
        """

        try:
            headers = {
                "Authorization": f"Bearer {self.api_key}",
                "Content-Type": "application/json"
            }
            payload = {
                "model": self.model,
                "messages": [
                    {"role": "system", "content": "You are a python-pptx expert. Return only valid Python code."},
                    {"role": "user", "content": prompt}
                ],
                "temperature": 0.3,
                "max_tokens": 800
            }

            response = await self.client.post(self.base_url, headers=headers, json=payload)
            response.raise_for_status()

            result = response.json()
            content_text = result["choices"][0]["message"].get("content", "").strip()

            # Extract code block
            if '```python' in content_text:
                content_text = content_text.split('```python')[1].split('```')[0].strip()
            elif '```' in content_text:
                content_text = content_text.split('```')[1].split('```')[0].strip()

            # Basic validation: must contain 'def layout_slide'
            if 'def layout_slide' in content_text:
                return content_text

        except Exception as e:
            logger.warning(f"Slide layout code generation failed: {e}")

        return self._generate_fallback_slide_code(slide_content)

    @staticmethod
    def _generate_fallback_slide_code(slide_content: Dict[str, Any]) -> str:
        """Generate a simple fallback layout function."""
        return '''def layout_slide(slide, colors, typography):
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    # Simple two-column layout
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_box.text_frame.paragraphs[0].text = slide_content.get("title", "")
    title_box.text_frame.paragraphs[0].font.size = Pt(28)
    title_box.text_frame.paragraphs[0].font.bold = True
'''

    @staticmethod
    def _get_fallback_typography(industry: str) -> Dict[str, Dict[str, Any]]:
        """Get industry-appropriate fallback typography."""
        font_map = {
            'technology': 'Segoe UI',
            'healthcare': 'Arial',
            'finance': 'Times New Roman',
            'security': 'Arial',
            'education': 'Verdana',
            'consulting': 'Garamond',
            'startup': 'Montserrat',
            'minimal': 'Helvetica',
        }
        font = font_map.get(industry.lower(), 'Calibri')
        body_font = 'Open Sans' if industry.lower() == 'startup' else font

        return {
            "title_slide": {"font_name": font, "font_size": 36, "bold": True, "color": "primary"},
            "slide_title": {"font_name": font, "font_size": 24, "bold": True, "color": "primary"},
            "body_text": {"font_name": body_font, "font_size": 14, "bold": False, "color": "text"},
            "caption": {"font_name": body_font, "font_size": 11, "bold": False, "color": "secondary"}
        }

    async def generate_complete_design_system(
        self,
        json_data: Dict[str, Any],
        company_info: Dict[str, Any]
    ) -> Dict[str, Any]:
        """Generate complete design system for the presentation.

        Args:
            json_data: Presentation content
            company_info: Company information

        Returns:
            Complete design system including colors, typography, and layouts
        """
        if self.use_mock:
            return await self.mock_generator.generate_complete_design_system(json_data, company_info)

        try:
            # Step 1: Analyze content for style
            logger.info("Analyzing content for design style...")
            style_analysis = await self.analyze_content_style(json_data, company_info)

            # Step 2: Generate color scheme
            logger.info("Generating color scheme...")
            colors = await self.generate_color_scheme(style_analysis)

            # Step 3: Generate typography using AI
            logger.info("Generating typography system...")
            typography = await self.generate_typography_system(style_analysis)

            design_system = {
                "style_analysis": style_analysis,
                "colors": colors,
                "typography": typography,
                "metadata": {
                    "generated_at": datetime.now().isoformat(),
                    "ai_model": self.model,
                    "version": "2.0"
                }
            }

            logger.info(f"Generated design system: {style_analysis.get('style')} style for {style_analysis.get('industry')} industry")
            return design_system

        except Exception as e:
            logger.error(f"Error generating design system: {e}")
            return {
                "style_analysis": {
                    "industry": "government",
                    "formality": "formal",
                    "style": "corporate",
                    "color_preference": "professional blue",
                    "audience": "executives"
                },
                "colors": {
                    "primary": "#2E75B6",
                    "secondary": "#8C8C8C",
                    "background": "#FFFFFF",
                    "text": "#1F1F1F",
                    "accent": "#00B050",
                    "light": "#F8F9FA"
                },
                "typography": self._get_fallback_typography("corporate"),
                "metadata": {
                    "generated_at": datetime.now().isoformat(),
                    "ai_model": "fallback",
                    "version": "2.0"
                }
            }


def create_ai_design_generator(api_key: str) -> Optional[AIDesignGenerator]:
    """
    Create AI design generator if API key is available.

    Args:
        api_key: Hugging Face API key

    Returns:
        AIDesignGenerator instance or None if no API key
    """
    if not api_key:
        logger.warning("No Hugging Face API key provided - AI design features disabled")
        return None

    try:
        return AIDesignGenerator(api_key)
    except Exception as e:
        logger.error(f"Failed to create AI design generator: {e}")
        return None