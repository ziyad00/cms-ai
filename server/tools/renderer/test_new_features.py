#!/usr/bin/env python3
"""Tests for all newly ported olama features:
- Minimal theme, get_theme_by_style, get_smart_theme, DesignSystemBuilder
- Detailed validate_design_system
- AI typography/layout generation
- Table rendering, image insertion, dynamic positioning, slide numbering
- Visual generator (ChartGenerator, DiagramGenerator)
- Content injector (ContentValidator, TenderDataMapper)
- Proposal layouts (SlideType, ProposalLayouts, ProposalTemplate)
"""

import os
import sys
import tempfile
import unittest

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from design_templates import (
    DesignTemplateLibrary, DesignSystemBuilder, BackgroundDesign, BackgroundType,
    get_design_system_for_content, validate_design_system,
)
from ai_design_generator import AIDesignGenerator
from render_pptx import (
    SmartDesignRules, SmartLayoutDetector, AIEnhancedPPTXRenderer,
)
from content_injector import ContentValidator, TenderDataMapper, prepare_proposal_content
from proposal_layouts import (
    SlideType, SlideLayout, ProposalLayouts, ProposalTemplate,
    STANDARD_PROPOSAL, QUICK_PITCH, TECHNICAL_PROPOSAL,
)
from visual_generator import ChartGenerator, DiagramGenerator
from json_processor import validate_json_structure, ProposalTemplateProcessor
from abstractions import (
    LayoutType, SlideContext, IBackgroundRenderer, ILayoutGenerator,
    IContentAnalyzer, IThemeProvider, IDesignSystemGenerator, IPresentationGenerator,
    RendererRegistry, BaseSlideGenerator, BasePresentationGenerator,
    PluginManager, GeneratorFactory, SystemConfig, ConfigManager,
)

from pptx import Presentation
from pptx.util import Inches


# === Minimal Theme Tests ===

class TestMinimalTheme(unittest.TestCase):
    """Tests for the new Minimal Clean theme."""

    def test_minimal_theme_exists(self):
        theme = DesignTemplateLibrary.MINIMAL_THEME
        self.assertEqual(theme.name, "Minimal Clean")

    def test_minimal_theme_charcoal_primary(self):
        theme = DesignTemplateLibrary.MINIMAL_THEME
        self.assertEqual(theme.colors['primary'], '#333333')

    def test_minimal_theme_red_accent(self):
        theme = DesignTemplateLibrary.MINIMAL_THEME
        self.assertEqual(theme.colors['accent'], '#E53E3E')

    def test_minimal_theme_helvetica_font(self):
        theme = DesignTemplateLibrary.MINIMAL_THEME
        self.assertEqual(theme.typography['title_slide']['font_name'], 'Helvetica')

    def test_minimal_theme_solid_background(self):
        theme = DesignTemplateLibrary.MINIMAL_THEME
        self.assertEqual(theme.background_design.type, BackgroundType.SOLID)

    def test_minimal_in_get_all_themes(self):
        themes = DesignTemplateLibrary.get_all_themes()
        self.assertIn('minimal', themes)
        self.assertEqual(len(themes), 10)

    def test_minimal_in_get_theme_by_name(self):
        theme = DesignTemplateLibrary.get_theme_by_name("Minimal Clean")
        self.assertIsNotNone(theme)
        self.assertEqual(theme.name, "Minimal Clean")


# === Font Fix Tests ===

class TestThemeFontFixes(unittest.TestCase):
    """Tests for updated theme fonts."""

    def test_startup_uses_montserrat(self):
        theme = DesignTemplateLibrary.STARTUP_THEME
        self.assertEqual(theme.typography['title_slide']['font_name'], 'Montserrat')
        self.assertEqual(theme.typography['slide_title']['font_name'], 'Montserrat')

    def test_startup_body_uses_open_sans(self):
        theme = DesignTemplateLibrary.STARTUP_THEME
        self.assertEqual(theme.typography['body_text']['font_name'], 'Open Sans')

    def test_consulting_uses_garamond_titles(self):
        theme = DesignTemplateLibrary.CONSULTING_THEME
        self.assertEqual(theme.typography['title_slide']['font_name'], 'Garamond')
        self.assertEqual(theme.typography['slide_title']['font_name'], 'Garamond')

    def test_consulting_body_uses_calibri(self):
        theme = DesignTemplateLibrary.CONSULTING_THEME
        self.assertEqual(theme.typography['body_text']['font_name'], 'Calibri')


# === get_theme_by_style Tests ===

class TestGetThemeByStyle(unittest.TestCase):
    """Tests for style-based theme lookup."""

    def test_modern_style(self):
        theme = DesignTemplateLibrary.get_theme_by_style('modern')
        self.assertEqual(theme.name, 'Modern Tech')

    def test_minimal_style(self):
        theme = DesignTemplateLibrary.get_theme_by_style('minimal')
        self.assertEqual(theme.name, 'Minimal Clean')

    def test_bold_style(self):
        theme = DesignTemplateLibrary.get_theme_by_style('bold')
        self.assertEqual(theme.name, 'Startup Dynamic')

    def test_elegant_style(self):
        theme = DesignTemplateLibrary.get_theme_by_style('elegant')
        self.assertEqual(theme.name, 'Consulting Executive')

    def test_formal_style(self):
        theme = DesignTemplateLibrary.get_theme_by_style('formal')
        self.assertEqual(theme.name, 'Government Official')

    def test_friendly_style(self):
        theme = DesignTemplateLibrary.get_theme_by_style('friendly')
        self.assertEqual(theme.name, 'Educational Friendly')

    def test_dark_style(self):
        theme = DesignTemplateLibrary.get_theme_by_style('dark')
        self.assertEqual(theme.name, 'Cybersecurity')

    def test_unknown_defaults_to_corporate(self):
        theme = DesignTemplateLibrary.get_theme_by_style('xyzabc')
        self.assertEqual(theme.name, 'Corporate Professional')

    def test_partial_match(self):
        theme = DesignTemplateLibrary.get_theme_by_style('minimalist design')
        self.assertEqual(theme.name, 'Minimal Clean')


# === get_smart_theme Tests ===

class TestGetSmartTheme(unittest.TestCase):
    """Tests for multi-factor theme selection."""

    def test_industry_only(self):
        theme = DesignTemplateLibrary.get_smart_theme(industry='healthcare')
        self.assertEqual(theme.name, 'Healthcare Professional')

    def test_style_only(self):
        theme = DesignTemplateLibrary.get_smart_theme(style='bold')
        self.assertEqual(theme.name, 'Startup Dynamic')

    def test_industry_overrides_style(self):
        # Industry has weight 3.0, style 2.0 — industry wins
        theme = DesignTemplateLibrary.get_smart_theme(industry='finance', style='modern')
        self.assertEqual(theme.name, 'Financial Services')

    def test_formality_formal(self):
        theme = DesignTemplateLibrary.get_smart_theme(formality='formal')
        # Should pick from formal set (corporate, finance, government)
        self.assertIn(theme.name, ['Corporate Professional', 'Financial Services', 'Government Official'])

    def test_audience_executive(self):
        theme = DesignTemplateLibrary.get_smart_theme(audience='c-suite executives')
        self.assertIn(theme.name, ['Consulting Executive', 'Corporate Professional'])

    def test_all_empty_defaults_to_corporate(self):
        theme = DesignTemplateLibrary.get_smart_theme()
        self.assertEqual(theme.name, 'Corporate Professional')

    def test_combined_factors(self):
        theme = DesignTemplateLibrary.get_smart_theme(
            industry='technology', style='modern', formality='business', audience='developers'
        )
        self.assertEqual(theme.name, 'Modern Tech')


# === DesignSystemBuilder Tests ===

class TestDesignSystemBuilder(unittest.TestCase):
    """Tests for the DesignSystemBuilder class."""

    def test_build_from_theme(self):
        builder = DesignSystemBuilder()
        theme = builder.from_theme(DesignTemplateLibrary.CORPORATE_THEME).build()
        self.assertEqual(theme.name, 'Corporate Professional')
        self.assertEqual(theme.colors['primary'], '#2E75B6')

    def test_override_colors(self):
        builder = DesignSystemBuilder()
        theme = (builder.from_theme(DesignTemplateLibrary.CORPORATE_THEME)
                 .with_colors({'primary': '#FF0000'})
                 .build())
        self.assertEqual(theme.colors['primary'], '#FF0000')
        # Other colors preserved
        self.assertEqual(theme.colors['text'], '#2C3E50')

    def test_override_typography(self):
        builder = DesignSystemBuilder()
        theme = (builder.from_theme(DesignTemplateLibrary.CORPORATE_THEME)
                 .with_typography({'title_slide': {'font_name': 'Arial', 'font_size': 40, 'bold': True, 'color': 'primary'}})
                 .build())
        self.assertEqual(theme.typography['title_slide']['font_name'], 'Arial')

    def test_with_name(self):
        builder = DesignSystemBuilder()
        theme = builder.from_theme(DesignTemplateLibrary.CORPORATE_THEME).with_name("Custom Theme").build()
        self.assertEqual(theme.name, "Custom Theme")

    def test_build_system_returns_dict(self):
        builder = DesignSystemBuilder()
        system = builder.from_theme(DesignTemplateLibrary.HEALTHCARE_THEME).build_system()
        self.assertIn('theme', system)
        self.assertIn('colors', system)
        self.assertIn('typography', system)
        self.assertEqual(system['colors']['primary'], '#48BB78')

    def test_build_without_from_theme_uses_defaults(self):
        builder = DesignSystemBuilder()
        theme = builder.build()
        self.assertEqual(theme.name, 'Custom')
        # Should use corporate defaults
        self.assertIn('primary', theme.colors)

    def test_with_background(self):
        bg = BackgroundDesign(type=BackgroundType.GRADIENT, primary_color='#000000', secondary_color='#FFFFFF')
        builder = DesignSystemBuilder()
        theme = builder.from_theme(DesignTemplateLibrary.MINIMAL_THEME).with_background(bg).build()
        self.assertEqual(theme.background_design.type, BackgroundType.GRADIENT)


# === Detailed validate_design_system Tests ===

class TestValidateDesignSystemDetailed(unittest.TestCase):
    """Tests for the enhanced validate_design_system with detailed errors."""

    def test_valid_system_returns_true(self):
        system = {
            'colors': {'primary': '#000000', 'text': '#333333'},
            'typography': {'title_slide': {'font_size': 36}, 'body_text': {'font_size': 14}},
        }
        self.assertTrue(validate_design_system(system))

    def test_missing_colors_returns_false(self):
        system = {'typography': {'title_slide': {}}}
        self.assertFalse(validate_design_system(system))

    def test_missing_typography_returns_false(self):
        system = {'colors': {'primary': '#000000'}}
        self.assertFalse(validate_design_system(system))

    def test_detailed_returns_list(self):
        system = {'colors': {}}
        errors = validate_design_system(system, detailed=True)
        self.assertIsInstance(errors, list)
        self.assertTrue(len(errors) > 0)

    def test_detailed_missing_required_color(self):
        system = {
            'colors': {'secondary': '#666666'},
            'typography': {'title_slide': {}, 'body_text': {}},
        }
        errors = validate_design_system(system, detailed=True)
        primary_errors = [e for e in errors if "'primary'" in e]
        self.assertTrue(len(primary_errors) > 0)

    def test_detailed_invalid_hex_color(self):
        system = {
            'colors': {'primary': 'red', 'text': '#333333'},
            'typography': {'title_slide': {}, 'body_text': {}},
        }
        errors = validate_design_system(system, detailed=True)
        hex_errors = [e for e in errors if 'Invalid hex' in e]
        self.assertTrue(len(hex_errors) > 0)

    def test_detailed_font_size_out_of_range(self):
        system = {
            'colors': {'primary': '#000000', 'text': '#333333'},
            'typography': {'title_slide': {'font_size': 100}, 'body_text': {'font_size': 14}},
        }
        errors = validate_design_system(system, detailed=True)
        size_errors = [e for e in errors if 'Font size' in e]
        self.assertTrue(len(size_errors) > 0)

    def test_detailed_valid_system_returns_empty_list(self):
        system = {
            'colors': {'primary': '#000000', 'text': '#333333'},
            'typography': {'title_slide': {'font_size': 36}, 'body_text': {'font_size': 14}},
        }
        errors = validate_design_system(system, detailed=True)
        self.assertEqual(errors, [])


# === AI Typography / Layout Code Tests ===

class TestAITypographyGeneration(unittest.TestCase):
    """Tests for AI typography and layout code generation (mock mode)."""

    def test_fallback_typography_corporate(self):
        result = AIDesignGenerator._get_fallback_typography('corporate')
        self.assertEqual(result['title_slide']['font_name'], 'Calibri')

    def test_fallback_typography_technology(self):
        result = AIDesignGenerator._get_fallback_typography('technology')
        self.assertEqual(result['title_slide']['font_name'], 'Segoe UI')

    def test_fallback_typography_healthcare(self):
        result = AIDesignGenerator._get_fallback_typography('healthcare')
        self.assertEqual(result['title_slide']['font_name'], 'Arial')

    def test_fallback_typography_finance(self):
        result = AIDesignGenerator._get_fallback_typography('finance')
        self.assertEqual(result['title_slide']['font_name'], 'Times New Roman')

    def test_fallback_typography_startup(self):
        result = AIDesignGenerator._get_fallback_typography('startup')
        self.assertEqual(result['title_slide']['font_name'], 'Montserrat')
        self.assertEqual(result['body_text']['font_name'], 'Open Sans')

    def test_fallback_typography_minimal(self):
        result = AIDesignGenerator._get_fallback_typography('minimal')
        self.assertEqual(result['title_slide']['font_name'], 'Helvetica')

    def test_fallback_slide_code_contains_function(self):
        code = AIDesignGenerator._generate_fallback_slide_code({'title': 'Test', 'content': []})
        self.assertIn('def layout_slide', code)
        self.assertIn('Inches', code)


# === Dynamic Positioning Tests ===

class TestDynamicPositioning(unittest.TestCase):
    """Tests for SmartDesignRules.calculate_dynamic_positioning."""

    def test_basic_positioning(self):
        pos = SmartDesignRules.calculate_dynamic_positioning(
            ["item1", "item2"], 10.0, 7.5, has_title=True
        )
        self.assertIn('x', pos)
        self.assertIn('y', pos)
        self.assertIn('width', pos)
        self.assertIn('height', pos)
        self.assertGreater(pos['width'], 0)
        self.assertGreater(pos['height'], 0)

    def test_dense_content_uses_full_width(self):
        items = ["x" * 100 for _ in range(10)]
        pos = SmartDesignRules.calculate_dynamic_positioning(items, 10.0, 7.5)
        self.assertAlmostEqual(pos['x'], 0.5)
        self.assertAlmostEqual(pos['width'], 9.0)

    def test_sparse_content_centered(self):
        pos = SmartDesignRules.calculate_dynamic_positioning(["short"], 10.0, 7.5)
        self.assertGreater(pos['x'], 1.0)
        self.assertLess(pos['width'], 8.0)

    def test_no_title_starts_higher(self):
        with_title = SmartDesignRules.calculate_dynamic_positioning(["a"], 10.0, 7.5, has_title=True)
        no_title = SmartDesignRules.calculate_dynamic_positioning(["a"], 10.0, 7.5, has_title=False)
        self.assertLess(no_title['y'], with_title['y'])


# === Table Layout Detection Tests ===

class TestTableLayoutDetection(unittest.TestCase):
    """Tests for table layout detection and rendering."""

    def test_detect_table_with_pipes(self):
        items = ["Name | Role | Team", "Alice | Dev | Backend", "Bob | PM | Product"]
        layout = SmartLayoutDetector.detect_layout("Staff", items)
        self.assertEqual(layout, "table")

    def test_detect_table_with_tabs(self):
        items = ["Name\tRole", "Alice\tDev"]
        layout = SmartLayoutDetector.detect_layout("Staff", items)
        self.assertEqual(layout, "table")

    def test_no_table_without_delimiters(self):
        items = ["Alice is a developer", "Bob is a PM"]
        layout = SmartLayoutDetector.detect_layout("Staff", items)
        self.assertNotEqual(layout, "table")


# === Table Rendering Tests ===

class TestTableRendering(unittest.TestCase):
    """Tests for _create_table and _render_table_layout."""

    def test_render_table_layout_produces_table(self):
        spec = {
            "layouts": [{
                "name": "staff",
                "placeholders": [
                    {"id": "title", "type": "text", "content": "Team Members"},
                    {"id": "body", "type": "text",
                     "content": "Name | Role | Team\nAlice | Dev | Backend\nBob | PM | Product"}
                ]
            }]
        }
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            output = f.name
        try:
            renderer = AIEnhancedPPTXRenderer()
            renderer.render_pptx_sync(spec, output)
            prs = Presentation(output)
            self.assertEqual(len(prs.slides), 1)
            slide = prs.slides[0]
            table_shapes = [s for s in slide.shapes if s.has_table]
            self.assertGreaterEqual(len(table_shapes), 1, "Should have a table")
        finally:
            os.unlink(output)


# === Slide Number Tests ===

class TestSlideNumbers(unittest.TestCase):
    """Tests for slide numbering."""

    def test_slides_have_numbers(self):
        spec = {
            "layouts": [
                {"name": "s1", "placeholders": [
                    {"id": "title", "type": "text", "content": "Slide 1"},
                    {"id": "body", "type": "text", "content": "Content A"}
                ]},
                {"name": "s2", "placeholders": [
                    {"id": "title", "type": "text", "content": "Slide 2"},
                    {"id": "body", "type": "text", "content": "Content B"}
                ]},
            ]
        }
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            output = f.name
        try:
            renderer = AIEnhancedPPTXRenderer()
            renderer.render_pptx_sync(spec, output)
            prs = Presentation(output)
            self.assertEqual(len(prs.slides), 2)
            # Check that slide number text exists
            slide1_texts = [s.text_frame.text for s in prs.slides[0].shapes if s.has_text_frame]
            self.assertTrue(any('1 / 2' in t for t in slide1_texts), "Slide 1 should have '1 / 2'")
            slide2_texts = [s.text_frame.text for s in prs.slides[1].shapes if s.has_text_frame]
            self.assertTrue(any('2 / 2' in t for t in slide2_texts), "Slide 2 should have '2 / 2'")
        finally:
            os.unlink(output)


# === ContentValidator Tests ===

class TestContentValidator(unittest.TestCase):
    """Tests for ContentValidator."""

    def test_validate_valid_content(self):
        content = {'slides': [{'title': 'Test', 'content': ['Item 1']}]}
        is_valid, errors = ContentValidator.validate_content(content)
        self.assertTrue(is_valid)
        self.assertEqual(errors, [])

    def test_validate_empty_content(self):
        is_valid, errors = ContentValidator.validate_content({})
        self.assertFalse(is_valid)

    def test_validate_no_slides(self):
        is_valid, errors = ContentValidator.validate_content({'slides': []})
        self.assertFalse(is_valid)

    def test_validate_long_title(self):
        content = {'slides': [{'title': 'x' * 100, 'content': ['Short']}]}
        is_valid, errors = ContentValidator.validate_content(content)
        self.assertFalse(is_valid)
        self.assertTrue(any('title exceeds' in e for e in errors))

    def test_sanitize_text_truncates(self):
        long_text = "a" * 200
        result = ContentValidator.sanitize_text(long_text, 'title')
        self.assertLessEqual(len(result), ContentValidator.LIMITS['title'])
        self.assertTrue(result.endswith('...'))

    def test_sanitize_text_strips_whitespace(self):
        result = ContentValidator.sanitize_text('  hello  ', 'body')
        self.assertEqual(result, 'hello')

    def test_sanitize_empty(self):
        result = ContentValidator.sanitize_text('', 'body')
        self.assertEqual(result, '')

    def test_clean_content_preserves_structure(self):
        content = {
            'slides': [
                {'title': '  Test  ', 'content': ['  Item 1  ']}
            ]
        }
        cleaned = ContentValidator.clean_content(content)
        self.assertEqual(cleaned['slides'][0]['title'], 'Test')
        self.assertEqual(cleaned['slides'][0]['content'][0], 'Item 1')

    def test_validate_layouts_format(self):
        content = {
            'layouts': [{
                'placeholders': [
                    {'id': 'title', 'type': 'title', 'content': 'Test'},
                    {'id': 'body', 'type': 'body', 'content': 'Content'},
                ]
            }]
        }
        is_valid, errors = ContentValidator.validate_content(content)
        self.assertTrue(is_valid)


# === TenderDataMapper Tests ===

class TestTenderDataMapper(unittest.TestCase):
    """Tests for TenderDataMapper."""

    def test_map_basic_proposal(self):
        data = {
            'company': {'name': 'Acme Corp'},
            'project': {'name': 'New Website', 'subtitle': 'Q1 2025'},
        }
        slides = TenderDataMapper.map_proposal_to_slides(data)
        self.assertGreater(len(slides), 0)
        self.assertEqual(slides[0]['title'], 'New Website')

    def test_map_with_team(self):
        data = {
            'project': {'name': 'Test'},
            'team': [
                {'name': 'Alice', 'role': 'Lead'},
                {'name': 'Bob', 'role': 'Dev'},
            ],
        }
        slides = TenderDataMapper.map_proposal_to_slides(data)
        team_slide = [s for s in slides if s['title'] == 'Our Team']
        self.assertEqual(len(team_slide), 1)
        self.assertEqual(len(team_slide[0]['content']), 2)

    def test_map_with_timeline(self):
        data = {
            'project': {'name': 'Test'},
            'timeline': [
                {'name': 'Phase 1', 'duration': '2 weeks'},
                {'name': 'Phase 2', 'duration': '4 weeks'},
            ],
        }
        slides = TenderDataMapper.map_proposal_to_slides(data)
        timeline_slide = [s for s in slides if s['title'] == 'Project Timeline']
        self.assertEqual(len(timeline_slide), 1)

    def test_map_with_pricing(self):
        data = {
            'project': {'name': 'Test'},
            'pricing': {
                'total': '$50,000',
                'items': [{'name': 'Design', 'cost': '$15,000'}],
            },
        }
        slides = TenderDataMapper.map_proposal_to_slides(data)
        pricing_slide = [s for s in slides if s['title'] == 'Pricing']
        self.assertEqual(len(pricing_slide), 1)

    def test_slides_to_spec(self):
        slides = [
            {'title': 'Title', 'content': ['Item 1', 'Item 2'], 'layout_hint': 'title'},
        ]
        spec = TenderDataMapper.slides_to_spec(slides)
        self.assertIn('layouts', spec)
        self.assertEqual(len(spec['layouts']), 1)
        self.assertEqual(spec['layouts'][0]['placeholders'][0]['content'], 'Title')

    def test_prepare_proposal_content(self):
        data = {
            'company': {'name': 'Test Co'},
            'project': {'name': 'Test Project', 'summary': 'A brief summary'},
        }
        spec = prepare_proposal_content(data)
        self.assertIn('layouts', spec)
        self.assertGreater(len(spec['layouts']), 0)


# === ProposalLayouts Tests ===

class TestProposalLayouts(unittest.TestCase):
    """Tests for proposal layout definitions."""

    def test_slide_type_enum(self):
        self.assertEqual(SlideType.TITLE.value, 'title')
        self.assertEqual(SlideType.CLOSING.value, 'closing')

    def test_get_all_layouts_returns_10(self):
        layouts = ProposalLayouts.get_all_layouts()
        self.assertEqual(len(layouts), 10)

    def test_get_required_layouts(self):
        required = ProposalLayouts.get_required_layouts()
        self.assertGreater(len(required), 0)
        for layout in required:
            self.assertTrue(layout.required)

    def test_get_layout_by_type(self):
        layout = ProposalLayouts.get_layout_by_type(SlideType.TITLE)
        self.assertIsNotNone(layout)
        self.assertEqual(layout.slide_type, SlideType.TITLE)

    def test_slide_layout_to_spec_layout(self):
        layout = ProposalLayouts.TITLE
        result = layout.to_spec_layout({
            'title': 'My Project',
            'subtitle': 'Q1 2025',
            'company': 'Acme',
        })
        self.assertEqual(result['name'], 'Title Slide')
        self.assertIn('placeholders', result)
        titles = [p for p in result['placeholders'] if p['id'] == 'title']
        self.assertEqual(titles[0]['content'], 'My Project')

    def test_standard_proposal_template(self):
        self.assertEqual(STANDARD_PROPOSAL.name, 'Standard Proposal')
        self.assertEqual(len(STANDARD_PROPOSAL.slide_order), 10)
        self.assertEqual(STANDARD_PROPOSAL.slide_order[0], SlideType.TITLE)
        self.assertEqual(STANDARD_PROPOSAL.slide_order[-1], SlideType.CLOSING)

    def test_quick_pitch_template(self):
        self.assertEqual(len(QUICK_PITCH.slide_order), 5)

    def test_technical_proposal_template(self):
        self.assertIn(SlideType.METHODOLOGY, TECHNICAL_PROPOSAL.slide_order)

    def test_proposal_template_generate_spec(self):
        content = {
            'title': {'title': 'Project X', 'company': 'Acme'},
            'executive_summary': {'summary': 'Brief overview'},
            'closing': {'contact': 'john@acme.com'},
        }
        spec = QUICK_PITCH.generate_spec(content)
        self.assertIn('layouts', spec)
        self.assertEqual(len(spec['layouts']), 5)

    def test_proposal_template_with_branding(self):
        template = ProposalTemplate(
            name="Branded",
            description="Test",
            slide_order=[SlideType.TITLE],
            branding={'colors': {'primary': '#FF0000'}, 'company': {'name': 'Test'}},
        )
        spec = template.generate_spec({'title': {'title': 'Test'}})
        self.assertIn('tokens', spec)
        self.assertEqual(spec['tokens']['colors']['primary'], '#FF0000')


# === Visual Generator Tests ===

class TestChartGenerator(unittest.TestCase):
    """Tests for ChartGenerator (matplotlib-based)."""

    def setUp(self):
        self.gen = ChartGenerator()

    def test_available_property(self):
        # Should be True if matplotlib installed, False otherwise
        self.assertIsInstance(self.gen.available, bool)

    def test_pie_chart_generates_file(self):
        if not self.gen.available:
            self.skipTest("matplotlib not installed")
        path = self.gen.generate_pie_chart(
            ['A', 'B', 'C'], [30, 50, 20], 'Test Pie'
        )
        self.assertIsNotNone(path)
        self.assertTrue(os.path.exists(path))
        self.assertGreater(os.path.getsize(path), 0)
        os.unlink(path)

    def test_bar_chart_generates_file(self):
        if not self.gen.available:
            self.skipTest("matplotlib not installed")
        path = self.gen.generate_bar_chart(
            ['X', 'Y', 'Z'], [100, 200, 150], 'Test Bar'
        )
        self.assertIsNotNone(path)
        self.assertTrue(os.path.exists(path))
        os.unlink(path)

    def test_progress_chart_generates_file(self):
        if not self.gen.available:
            self.skipTest("matplotlib not installed")
        path = self.gen.generate_progress_chart(
            [('Task A', 75), ('Task B', 50)], 'Progress'
        )
        self.assertIsNotNone(path)
        os.unlink(path)

    def test_gantt_chart_generates_file(self):
        if not self.gen.available:
            self.skipTest("matplotlib not installed")
        tasks = [
            {'name': 'Design', 'start': 0, 'duration': 3},
            {'name': 'Build', 'start': 2, 'duration': 5},
        ]
        path = self.gen.generate_gantt_chart(tasks, 'Gantt')
        self.assertIsNotNone(path)
        os.unlink(path)

    def test_empty_data_returns_none(self):
        path = self.gen.generate_pie_chart([], [], '')
        self.assertIsNone(path)

    def test_custom_colors(self):
        gen = ChartGenerator(colors={'primary': '#FF0000', 'accent': '#00FF00', 'secondary': '#0000FF'})
        colors = gen._get_chart_colors(3)
        self.assertEqual(colors[0], '#FF0000')


class TestDiagramGenerator(unittest.TestCase):
    """Tests for DiagramGenerator (Pillow-based)."""

    def setUp(self):
        self.gen = DiagramGenerator()

    def test_available_property(self):
        self.assertIsInstance(self.gen.available, bool)

    def test_architecture_diagram_generates_file(self):
        if not self.gen.available:
            self.skipTest("Pillow not installed")
        layers = [
            {'name': 'Frontend'},
            {'name': 'API Gateway'},
            {'name': 'Backend Services'},
            {'name': 'Database'},
        ]
        path = self.gen.generate_architecture_diagram(layers, 'Architecture')
        self.assertIsNotNone(path)
        self.assertTrue(os.path.exists(path))
        os.unlink(path)

    def test_org_chart_generates_file(self):
        if not self.gen.available:
            self.skipTest("Pillow not installed")
        nodes = [
            {'name': 'CEO'},
            {'name': 'CTO'},
            {'name': 'CFO'},
            {'name': 'VP Eng'},
        ]
        path = self.gen.generate_org_chart(nodes)
        self.assertIsNotNone(path)
        os.unlink(path)

    def test_empty_layers_returns_none(self):
        path = self.gen.generate_architecture_diagram([], '')
        self.assertIsNone(path)

    def test_single_node_org_chart(self):
        if not self.gen.available:
            self.skipTest("Pillow not installed")
        path = self.gen.generate_org_chart([{'name': 'CEO'}])
        self.assertIsNotNone(path)
        os.unlink(path)


# === get_design_system_for_content with Smart Theme Tests ===

class TestSmartDesignSystemForContent(unittest.TestCase):
    """Tests that get_design_system_for_content uses get_smart_theme."""

    def test_uses_industry(self):
        system = get_design_system_for_content(
            "Test content", {'industry': 'healthcare'}
        )
        self.assertEqual(system['theme'].name, 'Healthcare Professional')

    def test_uses_style(self):
        system = get_design_system_for_content(
            "Test content", {'industry': '', 'style': 'minimal'}
        )
        self.assertEqual(system['theme'].name, 'Minimal Clean')

    def test_uses_formality(self):
        system = get_design_system_for_content(
            "Test content", {'formality': 'formal'}
        )
        # Should be one of the formal themes
        self.assertIn(system['theme'].name,
                       ['Corporate Professional', 'Financial Services', 'Government Official'])


# === Enhanced TenderDataMapper (map_tender_to_proposal) Tests ===

class TestTenderDataMapperFull(unittest.TestCase):
    """Tests for TenderDataMapper.map_tender_to_proposal with all _create_* methods."""

    def setUp(self):
        self.tender_data = {
            'title': 'Cloud Migration Project',
            'reference_number': 'RFP-2025-001',
            'government_entity': 'Ministry of IT',
            'duration_months': 12,
            'estimated_value': '$500,000',
        }
        self.company_data = {
            'name': 'TechCorp',
            'address': '123 Main St',
            'phone': '+1-555-0100',
            'email': 'info@techcorp.com',
            'core_technologies': 'cloud infrastructure and AI',
        }

    def test_map_tender_returns_all_sections(self):
        result = TenderDataMapper.map_tender_to_proposal(self.tender_data, self.company_data)
        expected_keys = [
            'title', 'executive_summary', 'technical_approach', 'timeline',
            'team', 'experience', 'specifications', 'implementation',
            'risk_mitigation', 'timeline_data', 'budget_data', 'architecture_components',
        ]
        for key in expected_keys:
            self.assertIn(key, result, f"Missing section: {key}")

    def test_title_content_has_tender_title(self):
        result = TenderDataMapper.map_tender_to_proposal(self.tender_data, self.company_data)
        self.assertEqual(result['title']['tender_title'], 'Cloud Migration Project')

    def test_title_content_has_company_info(self):
        result = TenderDataMapper.map_tender_to_proposal(self.tender_data, self.company_data)
        self.assertIn('TechCorp', result['title']['company_info'])
        self.assertIn('123 Main St', result['title']['company_info'])

    def test_title_content_has_reference(self):
        result = TenderDataMapper.map_tender_to_proposal(self.tender_data, self.company_data)
        self.assertIn('RFP-2025-001', result['title']['tender_reference'])

    def test_executive_summary_has_overview(self):
        result = TenderDataMapper.map_tender_to_proposal(self.tender_data, self.company_data)
        summary = result['executive_summary']
        self.assertIn('overview', summary)
        self.assertGreater(len(summary['overview']), 0)
        self.assertIn('key_benefits', summary)
        self.assertIn('differentiators', summary)

    def test_technical_approach_has_stack(self):
        result = TenderDataMapper.map_tender_to_proposal(self.tender_data, self.company_data)
        tech = result['technical_approach']
        self.assertIn('methodology_overview', tech)
        self.assertIn('technical_stack', tech)
        self.assertGreater(len(tech['technical_stack']), 0)

    def test_timeline_has_phases(self):
        result = TenderDataMapper.map_tender_to_proposal(self.tender_data, self.company_data)
        timeline = result['timeline']
        self.assertIn('timeline_chart', timeline)
        self.assertGreater(len(timeline['timeline_chart']), 0)

    def test_team_has_table(self):
        result = TenderDataMapper.map_tender_to_proposal(self.tender_data, self.company_data)
        team = result['team']
        self.assertIn('team_table', team)
        self.assertGreater(len(team['team_table']), 0)
        # Each row should be pipe-delimited
        self.assertIn('|', team['team_table'][0])

    def test_team_uses_company_data_team(self):
        company = dict(self.company_data)
        company['team'] = [
            {"role": "CTO", "name": "Alice", "experience": "15 years", "qualifications": "PhD"},
        ]
        result = TenderDataMapper.map_tender_to_proposal(self.tender_data, company)
        self.assertIn('Alice', result['team']['team_table'][0])

    def test_experience_has_projects(self):
        result = TenderDataMapper.map_tender_to_proposal(self.tender_data, self.company_data)
        exp = result['experience']
        self.assertIn('project_1', exp)
        self.assertIn('achievements', exp)
        self.assertGreater(len(exp['achievements']), 0)

    def test_specifications_has_table(self):
        result = TenderDataMapper.map_tender_to_proposal(self.tender_data, self.company_data)
        specs = result['specifications']
        self.assertIn('specs_table', specs)
        self.assertGreater(len(specs['specs_table']), 0)
        self.assertIn('|', specs['specs_table'][0])

    def test_implementation_has_phases_and_deliverables(self):
        result = TenderDataMapper.map_tender_to_proposal(self.tender_data, self.company_data)
        impl = result['implementation']
        self.assertIn('phases', impl)
        self.assertIn('deliverables', impl)
        self.assertEqual(len(impl['phases']), 6)
        self.assertEqual(len(impl['deliverables']), 6)

    def test_risk_has_table(self):
        result = TenderDataMapper.map_tender_to_proposal(self.tender_data, self.company_data)
        risk = result['risk_mitigation']
        self.assertIn('risks_table', risk)
        self.assertEqual(len(risk['risks_table']), 4)

    def test_budget_included_when_present(self):
        tender = dict(self.tender_data)
        tender['budget_breakdown'] = {'Dev': 300000, 'Ops': 100000}
        result = TenderDataMapper.map_tender_to_proposal(tender, self.company_data)
        self.assertIn('budget', result)
        self.assertIn('budget_chart', result['budget'])
        self.assertIn('budget_breakdown', result['budget'])
        self.assertGreater(len(result['budget']['budget_breakdown']), 0)

    def test_budget_excluded_when_absent(self):
        result = TenderDataMapper.map_tender_to_proposal(self.tender_data, self.company_data)
        self.assertNotIn('budget', result)

    def test_timeline_data_for_visualization(self):
        result = TenderDataMapper.map_tender_to_proposal(self.tender_data, self.company_data)
        td = result['timeline_data']
        self.assertIsInstance(td, list)
        self.assertGreater(len(td), 0)
        self.assertIn('start_date', td[0])
        self.assertIn('end_date', td[0])

    def test_budget_data_for_visualization(self):
        result = TenderDataMapper.map_tender_to_proposal(self.tender_data, self.company_data)
        bd = result['budget_data']
        self.assertIsInstance(bd, dict)
        self.assertGreater(len(bd), 0)

    def test_architecture_components(self):
        result = TenderDataMapper.map_tender_to_proposal(self.tender_data, self.company_data)
        arch = result['architecture_components']
        self.assertIsInstance(arch, list)
        self.assertEqual(len(arch), 5)
        self.assertIn('name', arch[0])
        self.assertIn('x', arch[0])
        self.assertIn('y', arch[0])


# === validate_json_structure Tests ===

class TestValidateJsonStructure(unittest.TestCase):
    """Tests for validate_json_structure from json_processor."""

    def test_valid_structure(self):
        data = {'slides': [{'title': 'T', 'content': ['item'], 'slide_number': 1}]}
        errors = validate_json_structure(data)
        self.assertEqual(errors, [])

    def test_missing_slides_key(self):
        errors = validate_json_structure({'layouts': []})
        self.assertIn("Missing 'slides' key in JSON data", errors)

    def test_slides_not_a_list(self):
        errors = validate_json_structure({'slides': 'not a list'})
        self.assertIn("'slides' must be a list", errors)

    def test_empty_slides(self):
        errors = validate_json_structure({'slides': []})
        self.assertIn("No slides found in data", errors)

    def test_slide_not_a_dict(self):
        errors = validate_json_structure({'slides': ['not a dict']})
        self.assertTrue(any('Must be a dictionary' in e for e in errors))

    def test_missing_title(self):
        errors = validate_json_structure({'slides': [{'content': ['item']}]})
        self.assertTrue(any("Missing 'title'" in e for e in errors))

    def test_missing_content(self):
        errors = validate_json_structure({'slides': [{'title': 'T'}]})
        self.assertTrue(any("Missing 'content'" in e for e in errors))

    def test_content_not_a_list(self):
        errors = validate_json_structure({'slides': [{'title': 'T', 'content': 'not list'}]})
        self.assertTrue(any("'content' must be a list" in e for e in errors))

    def test_invalid_slide_number(self):
        errors = validate_json_structure({'slides': [{'title': 'T', 'content': [], 'slide_number': -1}]})
        self.assertTrue(any("positive integer" in e for e in errors))

    def test_non_dict_input(self):
        errors = validate_json_structure("not a dict")
        self.assertIn("JSON data must be a dictionary", errors)


# === ProposalTemplateProcessor Tests ===

class TestProposalTemplateProcessor(unittest.TestCase):
    """Tests for ProposalTemplateProcessor from json_processor."""

    def setUp(self):
        self.json_data = {
            'slides': [
                {'title': 'Proposal for [Company Name]', 'content': [
                    'Submitted by [Company Name]',
                    'Client: [Client]',
                    'Ref: [RFP Number]',
                    'Date: [Date]',
                ], 'slide_number': 1},
                {'title': 'Overview', 'content': ['Details here'], 'slide_number': 2},
            ]
        }

    def test_add_branding_replaces_company_name(self):
        result = ProposalTemplateProcessor.add_branding(
            self.json_data, {'name': 'Acme Inc', 'client': 'BigCorp', 'reference': 'RFP-42'}
        )
        self.assertIn('Acme Inc', result['slides'][0]['title'])
        self.assertIn('Acme Inc', result['slides'][0]['content'][0])

    def test_add_branding_replaces_client(self):
        result = ProposalTemplateProcessor.add_branding(
            self.json_data, {'name': 'Acme', 'client': 'BigCorp'}
        )
        self.assertIn('BigCorp', result['slides'][0]['content'][1])

    def test_add_branding_replaces_rfp_number(self):
        result = ProposalTemplateProcessor.add_branding(
            self.json_data, {'name': 'X', 'reference': 'RFP-42'}
        )
        self.assertIn('RFP-42', result['slides'][0]['content'][2])

    def test_add_branding_does_not_mutate_original(self):
        original_title = self.json_data['slides'][0]['title']
        ProposalTemplateProcessor.add_branding(self.json_data, {'name': 'Test'})
        self.assertEqual(self.json_data['slides'][0]['title'], original_title)

    def test_add_slide_numbers(self):
        data = {'slides': [{'title': 'A'}, {'title': 'B'}, {'title': 'C'}]}
        result = ProposalTemplateProcessor.add_slide_numbers(data)
        self.assertEqual(result['slides'][0]['slide_number'], 1)
        self.assertEqual(result['slides'][1]['slide_number'], 2)
        self.assertEqual(result['slides'][2]['slide_number'], 3)

    def test_add_slide_numbers_preserves_existing(self):
        data = {'slides': [{'title': 'A', 'slide_number': 10}, {'title': 'B'}]}
        result = ProposalTemplateProcessor.add_slide_numbers(data)
        self.assertEqual(result['slides'][0]['slide_number'], 10)
        self.assertEqual(result['slides'][1]['slide_number'], 2)

    def test_customize_content_replaces_title(self):
        result = ProposalTemplateProcessor.customize_content(
            self.json_data,
            {'slides': {'1': {'title': 'Custom Title'}}}
        )
        self.assertEqual(result['slides'][0]['title'], 'Custom Title')

    def test_customize_content_replaces_content(self):
        result = ProposalTemplateProcessor.customize_content(
            self.json_data,
            {'slides': {'2': {'content': ['New content']}}}
        )
        self.assertEqual(result['slides'][1]['content'], ['New content'])

    def test_customize_content_appends(self):
        result = ProposalTemplateProcessor.customize_content(
            self.json_data,
            {'slides': {'2': {'append_content': ['Extra item']}}}
        )
        self.assertIn('Extra item', result['slides'][1]['content'])
        self.assertIn('Details here', result['slides'][1]['content'])

    def test_customize_content_does_not_mutate_original(self):
        original_content = list(self.json_data['slides'][1]['content'])
        ProposalTemplateProcessor.customize_content(
            self.json_data,
            {'slides': {'2': {'content': ['Replaced']}}}
        )
        self.assertEqual(self.json_data['slides'][1]['content'], original_content)


# === Abstractions Tests ===

class TestLayoutTypeEnum(unittest.TestCase):
    """Tests for LayoutType enum."""

    def test_all_values(self):
        self.assertEqual(LayoutType.TITLE.value, 'title')
        self.assertEqual(LayoutType.CONTENT.value, 'content')
        self.assertEqual(LayoutType.COMPARISON.value, 'comparison')
        self.assertEqual(LayoutType.TIMELINE.value, 'timeline')
        self.assertEqual(LayoutType.METRICS.value, 'metrics')
        self.assertEqual(LayoutType.QUOTE.value, 'quote')
        self.assertEqual(LayoutType.DATA_VISUALIZATION.value, 'data_visualization')

    def test_enum_count(self):
        self.assertEqual(len(LayoutType), 7)


class TestSlideContext(unittest.TestCase):
    """Tests for SlideContext dataclass."""

    def test_creation_with_required_fields(self):
        ctx = SlideContext(
            slide_data={'title': 'Test'},
            design_system={'colors': {}},
            slide_number=1,
            total_slides=5,
        )
        self.assertEqual(ctx.slide_number, 1)
        self.assertEqual(ctx.total_slides, 5)
        self.assertIsNone(ctx.content_analysis)

    def test_creation_with_content_analysis(self):
        ctx = SlideContext(
            slide_data={}, design_system={}, slide_number=1, total_slides=1,
            content_analysis={'sentiment': 'positive'},
        )
        self.assertEqual(ctx.content_analysis['sentiment'], 'positive')


class TestRendererRegistry(unittest.TestCase):
    """Tests for RendererRegistry."""

    def test_empty_registry(self):
        reg = RendererRegistry()
        self.assertEqual(reg.background_renderer_count, 0)
        self.assertEqual(reg.layout_generator_count, 0)
        self.assertEqual(reg.content_analyzer_count, 0)
        self.assertEqual(reg.theme_provider_count, 0)

    def test_register_and_get_background_renderer(self):
        reg = RendererRegistry()

        class MockBgRenderer(IBackgroundRenderer):
            def render_background(self, slide, design_config):
                pass
            def supports_background_type(self, bg_type):
                return bg_type == 'gradient'

        renderer = MockBgRenderer()
        reg.register_background_renderer(renderer)
        self.assertEqual(reg.background_renderer_count, 1)
        self.assertIs(reg.get_background_renderer('gradient'), renderer)
        self.assertIsNone(reg.get_background_renderer('solid'))

    def test_register_layout_generator(self):
        reg = RendererRegistry()

        class MockLayoutGen(ILayoutGenerator):
            def generate_layout(self, context):
                return None
            def supports_layout_type(self, layout_type):
                return layout_type == LayoutType.TITLE

        gen = MockLayoutGen()
        reg.register_layout_generator(LayoutType.TITLE, gen)
        self.assertEqual(reg.layout_generator_count, 1)
        self.assertIs(reg.get_layout_generator(LayoutType.TITLE), gen)
        self.assertIsNone(reg.get_layout_generator(LayoutType.CONTENT))

    def test_register_content_analyzer(self):
        reg = RendererRegistry()

        class MockAnalyzer(IContentAnalyzer):
            def analyze_content(self, slide_data):
                return {}
            def detect_layout_type(self, content):
                return LayoutType.CONTENT

        analyzer = MockAnalyzer()
        reg.register_content_analyzer(analyzer)
        self.assertEqual(reg.content_analyzer_count, 1)
        self.assertIs(reg.get_content_analyzer(), analyzer)

    def test_register_theme_provider(self):
        reg = RendererRegistry()

        class MockTheme(IThemeProvider):
            def __init__(self):
                self.name = 'test'
            def get_theme_by_industry(self, industry):
                return {}
            def get_theme_by_style(self, style):
                return {}
            def list_available_themes(self):
                return ['test']

        provider = MockTheme()
        reg.register_theme_provider(provider)
        self.assertEqual(reg.theme_provider_count, 1)
        self.assertIs(reg.get_theme_provider(), provider)
        self.assertIs(reg.get_theme_provider('test'), provider)
        # When name doesn't match, falls back to first provider
        self.assertIs(reg.get_theme_provider('nonexistent'), provider)

    def test_get_content_analyzer_empty(self):
        reg = RendererRegistry()
        self.assertIsNone(reg.get_content_analyzer())

    def test_get_theme_provider_empty(self):
        reg = RendererRegistry()
        self.assertIsNone(reg.get_theme_provider())


class TestSystemConfig(unittest.TestCase):
    """Tests for SystemConfig and ConfigManager."""

    def test_default_config(self):
        cfg = SystemConfig()
        self.assertEqual(cfg.default_theme, 'corporate')
        self.assertEqual(cfg.default_generator, 'powerpoint')
        self.assertTrue(cfg.enable_ai_enhancement)
        self.assertFalse(cfg.cache_enabled)
        self.assertEqual(cfg.max_slides, 100)
        self.assertEqual(cfg.supported_formats, ['pptx', 'pdf'])

    def test_custom_config(self):
        cfg = SystemConfig(default_theme='minimal', max_slides=50)
        self.assertEqual(cfg.default_theme, 'minimal')
        self.assertEqual(cfg.max_slides, 50)

    def test_config_manager_get(self):
        mgr = ConfigManager()
        cfg = mgr.get_config()
        self.assertEqual(cfg.default_theme, 'corporate')

    def test_config_manager_update(self):
        mgr = ConfigManager()
        mgr.update_config(default_theme='healthcare', max_slides=20)
        self.assertEqual(mgr.config.default_theme, 'healthcare')
        self.assertEqual(mgr.config.max_slides, 20)

    def test_config_manager_ignores_unknown_keys(self):
        mgr = ConfigManager()
        mgr.update_config(nonexistent_key='value')
        self.assertFalse(hasattr(mgr.config, 'nonexistent_key'))


class TestPluginManager(unittest.TestCase):
    """Tests for PluginManager."""

    def test_creates_registry(self):
        pm = PluginManager()
        reg = pm.get_registry()
        self.assertIsInstance(reg, RendererRegistry)

    def test_load_default_plugins_does_not_error(self):
        pm = PluginManager()
        pm.load_default_plugins()  # Should not raise


class TestGeneratorFactory(unittest.TestCase):
    """Tests for GeneratorFactory."""

    def test_unknown_type_raises(self):
        with self.assertRaises(ValueError):
            GeneratorFactory.create_presentation_generator('unknown')

    def test_unknown_renderer_type_raises(self):
        with self.assertRaises(ValueError):
            GeneratorFactory.create_background_renderer('unknown')

    def test_register_and_create(self):
        class FakeGen(IPresentationGenerator):
            def __init__(self, registry):
                self.registry = registry
            async def generate_presentation(self, content, company_info, style_preference='auto'):
                return b''

        GeneratorFactory.register_generator('fake', FakeGen)
        gen = GeneratorFactory.create_presentation_generator('fake')
        self.assertIsInstance(gen, FakeGen)
        # Clean up
        del GeneratorFactory._generator_registry['fake']


# === Layout Hint from Spec Tests ===

class TestLayoutHintFromSpec(unittest.TestCase):
    """Tests that explicit layout_hint in spec.layout.name is used directly."""

    def test_known_layout_hint_used_directly(self):
        """When layout name matches a known layout, use it instead of auto-detecting."""
        spec = {
            "layouts": [{
                "name": "timeline",
                "placeholders": [
                    {"id": "title", "type": "text", "content": "Our Plan"},
                    {"id": "body", "type": "text", "content": "Step one\nStep two\nStep three"}
                ]
            }]
        }
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            output = f.name
        try:
            renderer = AIEnhancedPPTXRenderer()
            renderer.render_pptx_sync(spec, output)
            prs = Presentation(output)
            self.assertEqual(len(prs.slides), 1)
            # Timeline layout creates connector shapes (lines)
            slide = prs.slides[0]
            shape_types = [s.shape_type for s in slide.shapes if hasattr(s, 'shape_type')]
            # Should have rendered as timeline (has shapes), not simple text
            self.assertGreater(len(slide.shapes), 2, "Timeline should have multiple shapes")
        finally:
            os.unlink(output)

    def test_unknown_layout_name_falls_through_to_auto_detect(self):
        """When layout name is not a known layout, fall through to SmartLayoutDetector."""
        spec = {
            "layouts": [{
                "name": "My Custom Slide Title",
                "placeholders": [
                    {"id": "title", "type": "text", "content": "Metrics Overview"},
                    {"id": "body", "type": "text", "content": "Revenue: 50%\nGrowth: 30%\nRetention: 85%"}
                ]
            }]
        }
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            output = f.name
        try:
            renderer = AIEnhancedPPTXRenderer()
            renderer.render_pptx_sync(spec, output)
            prs = Presentation(output)
            self.assertEqual(len(prs.slides), 1)
            # Should still render (auto-detect picks metrics or simple)
        finally:
            os.unlink(output)

    def test_all_known_layouts_accepted(self):
        """All 10 known layout names should be recognized."""
        known = {"title", "quote", "timeline", "hierarchy", "comparison",
                 "metrics", "table", "grid", "multi_column", "simple"}
        # Build a spec with each known layout
        layouts = []
        for name in sorted(known):
            layouts.append({
                "name": name,
                "placeholders": [
                    {"id": "title", "type": "text", "content": f"Test {name}"},
                    {"id": "body", "type": "text", "content": "Item A\nItem B\nItem C\nItem D"}
                ]
            })
        spec = {"layouts": layouts}
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            output = f.name
        try:
            renderer = AIEnhancedPPTXRenderer()
            renderer.render_pptx_sync(spec, output)
            prs = Presentation(output)
            self.assertEqual(len(prs.slides), len(known))
        finally:
            os.unlink(output)

    def test_layout_hint_case_insensitive(self):
        """Layout name matching should be case-insensitive."""
        spec = {
            "layouts": [{
                "name": "Timeline",
                "placeholders": [
                    {"id": "title", "type": "text", "content": "Plan"},
                    {"id": "body", "type": "text", "content": "Phase 1\nPhase 2"}
                ]
            }]
        }
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            output = f.name
        try:
            renderer = AIEnhancedPPTXRenderer()
            renderer.render_pptx_sync(spec, output)
            prs = Presentation(output)
            self.assertEqual(len(prs.slides), 1)
        finally:
            os.unlink(output)


if __name__ == '__main__':
    unittest.main()
