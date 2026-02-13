#!/usr/bin/env python3
"""Tests for smart chart generation, layout detection, new themes, and olama-ported features."""
import os
import sys
import tempfile
import unittest

# Add renderer directory to path
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from render_pptx import (
    ContentAnalyzer, ContentType, SmartDesignRules, AdvancedLayoutEngine,
    DynamicChartGenerator, SmartLayoutDetector, AIEnhancedPPTXRenderer,
)
from design_templates import DesignTemplateLibrary
from abstract_background_renderer import BaseBackgroundRenderer

from pptx import Presentation
from pptx.util import Inches


# === ContentAnalyzer Tests ===

class TestContentAnalyzer(unittest.TestCase):
    """Tests for ContentAnalyzer (ported from olama)."""

    def setUp(self):
        self.analyzer = ContentAnalyzer()

    def test_detect_timeline_content_type(self):
        result = self.analyzer.analyze_content("Project Timeline", ["Phase 1", "Phase 2"])
        self.assertEqual(result["content_type"], ContentType.TIMELINE)

    def test_detect_data_driven_content_type(self):
        result = self.analyzer.analyze_content("Key Metrics", ["Revenue: $5M", "Growth: 30%"])
        self.assertEqual(result["content_type"], ContentType.DATA_DRIVEN)

    def test_detect_comparison_content_type(self):
        result = self.analyzer.analyze_content("Options Comparison", ["Option A", "Option B"])
        self.assertEqual(result["content_type"], ContentType.COMPARISON)

    def test_detect_hierarchy_content_type(self):
        result = self.analyzer.analyze_content("System Architecture", ["Frontend", "Backend", "DB"])
        self.assertEqual(result["content_type"], ContentType.HIERARCHY)

    def test_detect_quote_content_type(self):
        result = self.analyzer.analyze_content("Vision", ["Innovation drives everything we do"])
        self.assertEqual(result["content_type"], ContentType.QUOTE)

    def test_detect_list_items_for_many_items(self):
        result = self.analyzer.analyze_content("Features", ["A", "B", "C", "D"])
        self.assertEqual(result["content_type"], ContentType.LIST_ITEMS)

    def test_sentiment_urgent(self):
        result = self.analyzer.analyze_content("Alert", ["Critical security threat detected immediately"])
        self.assertEqual(result["sentiment"], "urgent")

    def test_sentiment_positive(self):
        result = self.analyzer.analyze_content("Results", ["Great success and growth achievement"])
        self.assertEqual(result["sentiment"], "positive")

    def test_sentiment_negative(self):
        result = self.analyzer.analyze_content("Issues", ["Multiple problems and failure in the system"])
        self.assertEqual(result["sentiment"], "negative")

    def test_sentiment_neutral(self):
        result = self.analyzer.analyze_content("Overview", ["Standard information about the project"])
        self.assertEqual(result["sentiment"], "neutral")

    def test_complexity_simple(self):
        result = self.analyzer.analyze_content("Agenda", ["Item one", "Item two"])
        self.assertEqual(result["complexity"], "simple")

    def test_complexity_complex(self):
        long_items = [
            "This is a very long and detailed description of something that takes many words to explain thoroughly and completely"
        ]
        result = self.analyzer.analyze_content("Details", long_items)
        self.assertEqual(result["complexity"], "complex")

    def test_has_numbers(self):
        result = self.analyzer.analyze_content("Stats", ["Revenue: $5M", "Users: 10000"])
        self.assertTrue(result["has_numbers"])

    def test_has_dates(self):
        result = self.analyzer.analyze_content("Schedule", ["Q1 2025: Launch", "Q2 2025: Scale"])
        self.assertTrue(result["has_dates"])

    def test_no_dates(self):
        result = self.analyzer.analyze_content("Info", ["Simple text with no dates"])
        self.assertFalse(result["has_dates"])

    def test_hierarchy_level_executive(self):
        result = self.analyzer.analyze_content("Executive Summary", ["Key point"])
        self.assertEqual(result["hierarchy_level"], 1)

    def test_hierarchy_level_section(self):
        result = self.analyzer.analyze_content("Introduction", ["Welcome"])
        self.assertEqual(result["hierarchy_level"], 2)

    def test_hierarchy_level_detail(self):
        result = self.analyzer.analyze_content("Implementation Details", ["Step 1"])
        self.assertEqual(result["hierarchy_level"], 3)

    def test_visual_weight_high_for_short_executive(self):
        result = self.analyzer.analyze_content("Executive Summary", ["Key point"])
        self.assertGreater(result["visual_weight"], 0.7)

    def test_visual_weight_low_for_long_complex(self):
        long_items = [f"Detail point {i} with extra words to increase count" for i in range(20)]
        result = self.analyzer.analyze_content("Detailed Analysis", long_items)
        self.assertLess(result["visual_weight"], 0.5)

    def test_visual_weight_clamped_0_1_to_1_0(self):
        result = self.analyzer.analyze_content("X", ["Y"])
        self.assertGreaterEqual(result["visual_weight"], 0.1)
        self.assertLessEqual(result["visual_weight"], 1.0)

    def test_key_concepts_extraction(self):
        result = self.analyzer.analyze_content(
            "AI Strategy", ["Machine Learning and NLP are key technologies for our API platform"]
        )
        self.assertIsInstance(result["key_concepts"], list)
        self.assertLessEqual(len(result["key_concepts"]), 5)

    def test_word_count(self):
        result = self.analyzer.analyze_content("Title", ["one two three"])
        self.assertEqual(result["word_count"], 4)  # "Title" + "one two three"


# === SmartDesignRules Tests ===

class TestSmartDesignRules(unittest.TestCase):
    """Tests for SmartDesignRules (ported from olama)."""

    def test_golden_ratio_spacing(self):
        spacing = SmartDesignRules.calculate_golden_ratio_spacing(20.0)
        self.assertAlmostEqual(spacing['tight'], 20.0 / 1.618, places=2)
        self.assertEqual(spacing['normal'], 20.0)
        self.assertAlmostEqual(spacing['loose'], 20.0 * 1.618, places=2)
        self.assertAlmostEqual(spacing['extra_loose'], 20.0 * 1.618 * 1.618, places=2)

    def test_font_sizes_title_range(self):
        sizes = SmartDesignRules.calculate_optimal_font_sizes(0.5, 50, ContentType.LIST_ITEMS)
        self.assertGreaterEqual(sizes['title'], 24)
        self.assertLessEqual(sizes['title'], 48)

    def test_font_sizes_body_range(self):
        sizes = SmartDesignRules.calculate_optimal_font_sizes(0.5, 50, ContentType.LIST_ITEMS)
        self.assertGreaterEqual(sizes['body'], 12)
        self.assertLessEqual(sizes['body'], 24)

    def test_font_sizes_caption_smaller_than_body(self):
        sizes = SmartDesignRules.calculate_optimal_font_sizes(0.5, 50, ContentType.LIST_ITEMS)
        self.assertLess(sizes['caption'], sizes['body'])

    def test_font_sizes_larger_for_high_visual_weight(self):
        low = SmartDesignRules.calculate_optimal_font_sizes(0.2, 50, ContentType.LIST_ITEMS)
        high = SmartDesignRules.calculate_optimal_font_sizes(0.9, 50, ContentType.LIST_ITEMS)
        self.assertGreater(high['title'], low['title'])
        self.assertGreater(high['body'], low['body'])

    def test_font_sizes_smaller_for_many_words(self):
        few = SmartDesignRules.calculate_optimal_font_sizes(0.5, 15, ContentType.LIST_ITEMS)
        many = SmartDesignRules.calculate_optimal_font_sizes(0.5, 100, ContentType.LIST_ITEMS)
        self.assertGreater(few['title'], many['title'])

    def test_font_sizes_quote_gets_larger_body(self):
        quote = SmartDesignRules.calculate_optimal_font_sizes(0.5, 20, ContentType.QUOTE)
        normal = SmartDesignRules.calculate_optimal_font_sizes(0.5, 20, ContentType.LIST_ITEMS)
        self.assertGreater(quote['body'], normal['body'])

    def test_font_sizes_data_driven_gets_smaller_body(self):
        data = SmartDesignRules.calculate_optimal_font_sizes(0.5, 50, ContentType.DATA_DRIVEN)
        normal = SmartDesignRules.calculate_optimal_font_sizes(0.5, 50, ContentType.LIST_ITEMS)
        self.assertLess(data['body'], normal['body'])

    def test_contrast_ratio_dark_text_on_light_bg(self):
        result = SmartDesignRules.ensure_contrast_ratio('#FFFFFF', '#333333')
        self.assertEqual(result, '#1F1F1F')

    def test_contrast_ratio_light_text_on_dark_bg(self):
        result = SmartDesignRules.ensure_contrast_ratio('#1A202C', '#FFFFFF')
        self.assertEqual(result, '#FFFFFF')

    def test_contrast_ratio_light_bg_starting_with_F(self):
        result = SmartDesignRules.ensure_contrast_ratio('#F8F9FA', '#333333')
        self.assertEqual(result, '#1F1F1F')


# === AdvancedLayoutEngine Tests ===

class TestAdvancedLayoutEngine(unittest.TestCase):
    """Tests for AdvancedLayoutEngine (ported from olama)."""

    def test_detect_comparison_pattern(self):
        items = ["Current vs Proposed", "Old approach", "New approach"]
        self.assertEqual(AdvancedLayoutEngine.detect_layout_pattern(items), "comparison_columns")

    def test_detect_timeline_pattern(self):
        items = ["Q1 2025: Plan", "Q2 2025: Build", "Q3 2025: Launch"]
        self.assertEqual(AdvancedLayoutEngine.detect_layout_pattern(items), "horizontal_timeline")

    def test_detect_metrics_grid(self):
        items = ["Revenue: 50%", "Growth: 30%", "Retention: 85%"]
        self.assertEqual(AdvancedLayoutEngine.detect_layout_pattern(items), "metrics_grid")

    def test_detect_two_column(self):
        items = ["Point A", "Point B"]
        self.assertEqual(AdvancedLayoutEngine.detect_layout_pattern(items), "two_column")

    def test_detect_three_column(self):
        items = ["A", "B", "C"]
        self.assertEqual(AdvancedLayoutEngine.detect_layout_pattern(items), "three_column")

    def test_detect_grid_layout(self):
        items = ["A", "B", "C", "D", "E"]
        self.assertEqual(AdvancedLayoutEngine.detect_layout_pattern(items), "grid_layout")

    def test_detect_multi_column(self):
        items = [f"Item {i}" for i in range(10)]
        self.assertEqual(AdvancedLayoutEngine.detect_layout_pattern(items), "multi_column")

    def test_detect_single_column_empty(self):
        self.assertEqual(AdvancedLayoutEngine.detect_layout_pattern([]), "single_column")

    def test_optimal_columns_1_for_few(self):
        result = AdvancedLayoutEngine.calculate_optimal_columns(2, 10.0)
        self.assertEqual(result["columns"], 1)

    def test_optimal_columns_2_for_medium(self):
        result = AdvancedLayoutEngine.calculate_optimal_columns(4, 10.0)
        self.assertEqual(result["columns"], 2)

    def test_optimal_columns_3_for_many(self):
        result = AdvancedLayoutEngine.calculate_optimal_columns(7, 10.0)
        self.assertEqual(result["columns"], 3)

    def test_optimal_columns_4_for_lots(self):
        result = AdvancedLayoutEngine.calculate_optimal_columns(12, 10.0)
        self.assertEqual(result["columns"], 4)

    def test_optimal_columns_has_spacing(self):
        result = AdvancedLayoutEngine.calculate_optimal_columns(4, 10.0)
        self.assertIn("spacing", result)
        self.assertIn("margin", result)
        self.assertGreater(result["spacing"], 0)


# === SmartLayoutDetector Tests (updated with new layouts) ===

class TestSmartLayoutDetector(unittest.TestCase):
    """Tests for SmartLayoutDetector layout detection."""

    def test_detect_quote_single_short_item(self):
        layout = SmartLayoutDetector.detect_layout(
            "Vision", ["Innovation drives everything we do"]
        )
        self.assertEqual(layout, "quote")

    def test_detect_quote_not_for_long_single_item(self):
        long_text = " ".join(["word"] * 35)
        layout = SmartLayoutDetector.detect_layout("Quote", [long_text])
        self.assertNotEqual(layout, "quote")

    def test_detect_hierarchy_from_title(self):
        layout = SmartLayoutDetector.detect_layout(
            "System Architecture", ["Frontend", "Backend", "Database"]
        )
        self.assertEqual(layout, "hierarchy")

    def test_detect_hierarchy_structure_keyword(self):
        layout = SmartLayoutDetector.detect_layout(
            "Organization Structure", ["CEO", "CTO", "VP Engineering"]
        )
        self.assertEqual(layout, "hierarchy")

    def test_detect_grid_for_4_items(self):
        layout = SmartLayoutDetector.detect_layout(
            "Features", ["Feature A", "Feature B", "Feature C", "Feature D"]
        )
        self.assertEqual(layout, "grid")

    def test_detect_grid_for_6_items(self):
        layout = SmartLayoutDetector.detect_layout(
            "Benefits", [f"Benefit {i}" for i in range(6)]
        )
        self.assertEqual(layout, "grid")

    def test_detect_timeline_from_title(self):
        layout = SmartLayoutDetector.detect_layout(
            "Project Timeline", ["Phase 1", "Phase 2", "Phase 3"]
        )
        self.assertEqual(layout, "timeline")

    def test_detect_timeline_from_content(self):
        layout = SmartLayoutDetector.detect_layout(
            "Product Roadmap", ["Q1 2025: Foundation", "Q2 2025: Launch"]
        )
        self.assertEqual(layout, "timeline")

    def test_detect_comparison_from_title(self):
        layout = SmartLayoutDetector.detect_layout(
            "Current vs Proposed", ["Item A", "Item B", "Item C", "Item D"]
        )
        self.assertEqual(layout, "comparison")

    def test_detect_comparison_from_content(self):
        layout = SmartLayoutDetector.detect_layout(
            "Options", ["Current system handles requests", "Proposed solution improves speed",
                         "Current has limitations", "Proposed removes them"]
        )
        self.assertEqual(layout, "comparison")

    def test_detect_metrics_from_title(self):
        layout = SmartLayoutDetector.detect_layout(
            "Key Performance Metrics", ["Revenue: $5M", "Growth: 30%"]
        )
        self.assertEqual(layout, "metrics")

    def test_detect_metrics_from_percentages(self):
        layout = SmartLayoutDetector.detect_layout(
            "Results", ["Metric A: 40%", "Metric B: 60%"]
        )
        self.assertEqual(layout, "metrics")

    def test_detect_multi_column_many_items(self):
        items = [f"Item {i}" for i in range(8)]
        layout = SmartLayoutDetector.detect_layout("Overview", items)
        self.assertEqual(layout, "multi_column")

    def test_detect_simple_for_short_content(self):
        layout = SmartLayoutDetector.detect_layout(
            "Introduction", ["Welcome to the presentation", "Today we will cover..."]
        )
        self.assertEqual(layout, "simple")

    def test_detect_simple_for_empty_content(self):
        layout = SmartLayoutDetector.detect_layout("Title Only", [])
        self.assertEqual(layout, "simple")


# === DynamicChartGenerator Tests ===

class TestDynamicChartGenerator(unittest.TestCase):
    """Tests for DynamicChartGenerator data pattern detection."""

    def test_detect_percentages_returns_pie(self):
        items = ["Market Share: 34%", "Customer Retention: 92%", "Revenue Growth: 127%"]
        result = DynamicChartGenerator.detect_data_pattern(items)
        self.assertEqual(result["chart_type"], "pie")
        self.assertTrue(result["has_percentages"])
        self.assertEqual(len(result["data"]), 3)
        self.assertAlmostEqual(result["data"][0], 34.0)

    def test_detect_single_percentage_no_chart(self):
        items = ["Overall satisfaction: 85%"]
        result = DynamicChartGenerator.detect_data_pattern(items)
        self.assertIsNone(result["chart_type"])
        self.assertTrue(result["has_percentages"])
        self.assertEqual(len(result["data"]), 1)

    def test_detect_timeline_returns_line(self):
        items = ["Q1 2025: $2M", "Q2 2025: $3M", "Q3 2025: $4.5M"]
        result = DynamicChartGenerator.detect_data_pattern(items)
        self.assertEqual(result["chart_type"], "line")
        self.assertTrue(result["has_timeline"])

    def test_detect_numeric_data_returns_bar(self):
        items = ["Product A: 150 units", "Product B: 230 units", "Product C: 80 units"]
        result = DynamicChartGenerator.detect_data_pattern(items)
        self.assertEqual(result["chart_type"], "bar")
        self.assertEqual(len(result["data"]), 3)

    def test_no_data_returns_none(self):
        items = ["Simple text only", "No numbers here"]
        result = DynamicChartGenerator.detect_data_pattern(items)
        self.assertIsNone(result["chart_type"])
        self.assertEqual(len(result["data"]), 0)

    def test_labels_truncated_to_30_chars(self):
        items = ["This is a very long label that should be truncated at thirty: 50%",
                 "Short: 25%"]
        result = DynamicChartGenerator.detect_data_pattern(items)
        for label in result["labels"]:
            self.assertLessEqual(len(label), 30)

    def test_create_chart_on_slide(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        chart_data = {
            "chart_type": "pie",
            "data": [30, 40, 30],
            "labels": ["A", "B", "C"],
            "has_percentages": True,
            "has_timeline": False,
        }
        chart = DynamicChartGenerator.create_chart(
            slide, chart_data,
            Inches(1), Inches(1), Inches(5), Inches(4)
        )
        self.assertIsNotNone(chart)
        chart_shapes = [s for s in slide.shapes if s.has_chart]
        self.assertEqual(len(chart_shapes), 1)

    def test_create_chart_bar_type(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        chart_data = {
            "chart_type": "bar",
            "data": [100, 200, 150],
            "labels": ["X", "Y", "Z"],
            "has_percentages": False,
            "has_timeline": False,
        }
        chart = DynamicChartGenerator.create_chart(
            slide, chart_data,
            Inches(1), Inches(1), Inches(5), Inches(4)
        )
        self.assertIsNotNone(chart)

    def test_create_chart_empty_data_returns_none(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        chart_data = {"chart_type": "pie", "data": [], "labels": [],
                      "has_percentages": False, "has_timeline": False}
        chart = DynamicChartGenerator.create_chart(
            slide, chart_data, Inches(1), Inches(1), Inches(5), Inches(4)
        )
        self.assertIsNone(chart)

    def test_create_progress_bar(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        colors = {"light": "#E5E5E5", "accent": "#00B050", "text": "#2C3E50"}
        initial_shapes = len(slide.shapes)

        DynamicChartGenerator.create_progress_bar(
            slide, Inches(1), Inches(3), Inches(8), Inches(0.5),
            75.0, "Completion", colors
        )
        # Should add 3 shapes: bg bar, progress bar, label
        self.assertEqual(len(slide.shapes) - initial_shapes, 3)


# === Theme Tests ===

class TestNewThemes(unittest.TestCase):
    """Tests for new themes (Startup, Government, Consulting)."""

    def test_startup_theme_exists(self):
        theme = DesignTemplateLibrary.get_theme_for_industry("startup")
        self.assertEqual(theme.name, "Startup Dynamic")
        self.assertEqual(theme.colors["background"], "#1A202C")

    def test_government_theme_exists(self):
        theme = DesignTemplateLibrary.get_theme_for_industry("government")
        self.assertEqual(theme.name, "Government Official")

    def test_consulting_theme_exists(self):
        theme = DesignTemplateLibrary.get_theme_for_industry("consulting")
        self.assertEqual(theme.name, "Consulting Executive")
        self.assertEqual(theme.colors["accent"], "#D69E2E")

    def test_venture_maps_to_startup(self):
        theme = DesignTemplateLibrary.get_theme_for_industry("venture capital")
        self.assertEqual(theme.name, "Startup Dynamic")

    def test_advisory_maps_to_consulting(self):
        theme = DesignTemplateLibrary.get_theme_for_industry("advisory firm")
        self.assertEqual(theme.name, "Consulting Executive")

    def test_municipal_maps_to_government(self):
        theme = DesignTemplateLibrary.get_theme_for_industry("municipal services")
        self.assertEqual(theme.name, "Government Official")

    def test_get_all_themes_returns_10(self):
        themes = DesignTemplateLibrary.get_all_themes()
        self.assertEqual(len(themes), 10)

    def test_healthcare_has_cross_decorative(self):
        theme = DesignTemplateLibrary.HEALTHCARE_THEME
        bg = theme.background_design
        self.assertIsNotNone(bg)
        cross_elements = [e for e in bg.decorative_elements if e.get("shape_type") == "cross"]
        self.assertEqual(len(cross_elements), 1)

    def test_financial_has_watermark(self):
        theme = DesignTemplateLibrary.FINANCIAL_THEME
        self.assertIsNotNone(theme.watermark)
        self.assertEqual(theme.watermark["content"], "CONFIDENTIAL")

    def test_financial_has_header_footer_bars(self):
        theme = DesignTemplateLibrary.FINANCIAL_THEME
        bg = theme.background_design
        rects = [e for e in bg.decorative_elements if e.get("shape_type") == "rectangle"]
        self.assertEqual(len(rects), 2)

    def test_startup_has_hexagon_background(self):
        theme = DesignTemplateLibrary.STARTUP_THEME
        bg = theme.background_design
        self.assertIsNotNone(bg)
        self.assertEqual(bg.type.value, "hexagon_grid")


# === Cross Shape Tests ===

class TestCrossShapeRendering(unittest.TestCase):
    """Test that cross shapes render correctly."""

    def test_cross_shape_adds_two_rectangles(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        renderer = BaseBackgroundRenderer.__new__(BaseBackgroundRenderer)
        initial = len(slide.shapes)

        renderer._add_cross_shape(
            slide, Inches(1), Inches(1), Inches(2), Inches(2), "#4299E1"
        )
        self.assertEqual(len(slide.shapes) - initial, 2)


# === End-to-End Renderer Tests ===

class TestRendererEndToEnd(unittest.TestCase):
    """End-to-end tests for the full renderer with all layout types."""

    def _render_spec(self, spec):
        """Helper to render a spec and return the presentation."""
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            output = f.name
        try:
            renderer = AIEnhancedPPTXRenderer()
            renderer.render_pptx_sync(spec, output)
            return Presentation(output), output
        except Exception:
            os.unlink(output)
            raise

    def test_render_with_charts_produces_valid_pptx(self):
        spec = {
            "tokens": {"colors": {"primary": "#2E75B6", "text": "#2C3E50",
                                   "background": "#FFFFFF", "accent": "#3498DB",
                                   "light": "#F8F9FA"}},
            "layouts": [{
                "name": "metrics",
                "placeholders": [
                    {"id": "title", "type": "text", "content": "Market Performance Metrics"},
                    {"id": "body", "type": "text",
                     "content": "Share: 34%\nRetention: 92%\nGrowth: 127%"}
                ]
            }]
        }
        prs, output = self._render_spec(spec)
        try:
            self.assertEqual(len(prs.slides), 1)
            slide = prs.slides[0]
            chart_shapes = [s for s in slide.shapes if s.has_chart]
            self.assertGreaterEqual(len(chart_shapes), 1, "Expected at least one chart")
        finally:
            os.unlink(output)

    def test_render_timeline_produces_cards(self):
        spec = {
            "layouts": [{
                "name": "roadmap",
                "placeholders": [
                    {"id": "title", "type": "text", "content": "Project Phases"},
                    {"id": "body", "type": "text",
                     "content": "Phase 1: Plan\nPhase 2: Build\nPhase 3: Launch"}
                ]
            }]
        }
        prs, output = self._render_spec(spec)
        try:
            slide = prs.slides[0]
            self.assertGreaterEqual(len(slide.shapes), 7)
        finally:
            os.unlink(output)

    def test_render_comparison_produces_two_columns(self):
        spec = {
            "layouts": [{
                "name": "compare",
                "placeholders": [
                    {"id": "title", "type": "text", "content": "Current vs Proposed"},
                    {"id": "body", "type": "text",
                     "content": "Current: Old\nCurrent: Slow\nProposed: New\nProposed: Fast"}
                ]
            }]
        }
        prs, output = self._render_spec(spec)
        try:
            slide = prs.slides[0]
            self.assertGreaterEqual(len(slide.shapes), 5)
        finally:
            os.unlink(output)

    def test_render_quote_layout(self):
        spec = {
            "layouts": [{
                "name": "vision",
                "placeholders": [
                    {"id": "title", "type": "text", "content": "Our Vision"},
                    {"id": "body", "type": "text",
                     "content": "Innovation drives everything we do"}
                ]
            }]
        }
        prs, output = self._render_spec(spec)
        try:
            slide = prs.slides[0]
            # Quote: title text box + quote text box + accent line = 3+ shapes (plus background)
            self.assertGreaterEqual(len(slide.shapes), 3)
            # Check that quote text exists
            texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
            quote_found = any('\u201C' in t for t in texts)
            self.assertTrue(quote_found, "Quote marks should be present")
        finally:
            os.unlink(output)

    def test_render_grid_layout(self):
        spec = {
            "layouts": [{
                "name": "features",
                "placeholders": [
                    {"id": "title", "type": "text", "content": "Key Features"},
                    {"id": "body", "type": "text",
                     "content": "Fast Performance\nEasy Setup\nSecure Access\nScalable Design"}
                ]
            }]
        }
        prs, output = self._render_spec(spec)
        try:
            slide = prs.slides[0]
            # Grid: title + 4 cards (rounded rect) + 4 text boxes = 9+ shapes
            self.assertGreaterEqual(len(slide.shapes), 9)
        finally:
            os.unlink(output)

    def test_render_hierarchy_layout(self):
        spec = {
            "layouts": [{
                "name": "arch",
                "placeholders": [
                    {"id": "title", "type": "text", "content": "System Architecture"},
                    {"id": "body", "type": "text",
                     "content": "Platform Core\nFrontend Layer\nBackend Services\nDatabase"}
                ]
            }]
        }
        prs, output = self._render_spec(spec)
        try:
            slide = prs.slides[0]
            # Hierarchy: title + top box + top text + 3*(connector + child box + child text) = 13+ shapes
            self.assertGreaterEqual(len(slide.shapes), 10)
        finally:
            os.unlink(output)

    def test_render_multi_slide_with_mixed_layouts(self):
        spec = {
            "tokens": {"colors": {"primary": "#2E75B6", "background": "#FFFFFF",
                                   "text": "#2C3E50", "accent": "#3498DB", "light": "#F8F9FA"}},
            "layouts": [
                {
                    "name": "title",
                    "placeholders": [
                        {"id": "title", "type": "text", "content": "Company Overview"},
                        {"id": "body", "type": "text", "content": "Building the future of AI"}
                    ]
                },
                {
                    "name": "features",
                    "placeholders": [
                        {"id": "title", "type": "text", "content": "Key Features"},
                        {"id": "body", "type": "text",
                         "content": "Speed\nReliability\nSecurity\nScalability"}
                    ]
                },
                {
                    "name": "metrics",
                    "placeholders": [
                        {"id": "title", "type": "text", "content": "Growth Metrics"},
                        {"id": "body", "type": "text",
                         "content": "Revenue: 45%\nUsers: 78%\nRetention: 92%"}
                    ]
                },
                {
                    "name": "arch",
                    "placeholders": [
                        {"id": "title", "type": "text", "content": "System Architecture"},
                        {"id": "body", "type": "text",
                         "content": "API Gateway\nMicroservices\nEvent Bus\nData Lake"}
                    ]
                },
            ]
        }
        prs, output = self._render_spec(spec)
        try:
            self.assertEqual(len(prs.slides), 4)
            # Each slide should have shapes
            for i, slide in enumerate(prs.slides):
                self.assertGreater(len(slide.shapes), 0, f"Slide {i} should have shapes")
        finally:
            os.unlink(output)


if __name__ == '__main__':
    unittest.main()
