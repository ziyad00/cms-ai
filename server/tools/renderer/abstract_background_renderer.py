"""Abstract background renderer implementations."""
from typing import Dict, Any, List
from abc import abstractmethod
import math

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR


class BackgroundType:
    """Background type constants."""
    SOLID = "solid"
    GRADIENT = "gradient"
    DARK_GRADIENT = "dark_gradient"
    DIAGONAL_LINES = "diagonal_lines"
    HEXAGON_GRID = "hexagon_grid"
    MEDICAL_CURVES = "medical_curves"
    TECH_CIRCUIT = "tech_circuit"
    CORPORATE_BARS = "corporate_bars"


class IBackgroundRenderer:
    """Interface for background renderers."""

    def supports_background_type(self, bg_type: str) -> bool:
        """Check if this renderer supports the background type."""
        pass

    def render_background(self, slide, design_config: Dict[str, Any]) -> None:
        """Render background on the given slide."""
        pass


class BaseBackgroundRenderer(IBackgroundRenderer):
    """Base implementation for background renderers."""

    def render_background(self, slide, design_config: Dict[str, Any]) -> None:
        """Render background on the given slide."""
        background_design = design_config.get('background_design')
        if not background_design:
            return

        # Apply base background
        self._apply_base_background(slide, background_design)

        # Apply patterns
        self._apply_patterns(slide, background_design)

        # Add decorative elements
        decorative_elements = getattr(background_design, 'decorative_elements', [])
        self._add_decorative_elements(slide, decorative_elements)

        # Add watermark
        watermark = design_config.get('watermark')
        if watermark:
            self._add_watermark(slide, watermark)

    @abstractmethod
    def _apply_patterns(self, slide, background_design) -> None:
        """Apply pattern-specific rendering."""
        pass

    @staticmethod
    def _normalize_bg_type(bg_type) -> str:
        """Normalize BackgroundType to string (handles both Enum and plain string)."""
        if hasattr(bg_type, 'value'):
            return bg_type.value
        return str(bg_type)

    def _apply_base_background(self, slide, background_design) -> None:
        """Apply base background color/gradient."""
        bg_type = self._normalize_bg_type(getattr(background_design, 'type', BackgroundType.SOLID))
        primary_color = getattr(background_design, 'primary_color', '#FFFFFF')
        secondary_color = getattr(background_design, 'secondary_color', '#E0E0E0')

        if bg_type in [BackgroundType.SOLID, BackgroundType.MEDICAL_CURVES]:
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = self._hex_to_rgb(primary_color)
        elif bg_type in [BackgroundType.GRADIENT, BackgroundType.DARK_GRADIENT]:
            slide.background.fill.gradient()
            slide.background.fill.gradient_angle = 45
            gradient_stops = slide.background.fill.gradient_stops
            gradient_stops[0].color.rgb = self._hex_to_rgb(primary_color)
            if secondary_color:
                gradient_stops[1].color.rgb = self._hex_to_rgb(secondary_color)
        else:
            # Default: apply solid background for all other types
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = self._hex_to_rgb(primary_color)

    def _add_decorative_elements(self, slide, elements: List[Dict[str, Any]]) -> None:
        """Add decorative elements to the slide."""
        for element in elements or []:
            self._add_decorative_element(slide, element)

    def _add_decorative_element(self, slide, element: Dict[str, Any]) -> None:
        """Add a single decorative element."""
        shape_type = element.get("shape_type", "rectangle")
        position = element.get("position", {})
        color = element.get("color", "#000000")

        x = Inches(position.get("x", 0))
        y = Inches(position.get("y", 0))
        width = Inches(position.get("width", 1))
        height = Inches(position.get("height", 1))

        if shape_type == "rectangle":
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = self._hex_to_rgb(color)
            shape.line.fill.background()
        elif shape_type == "circle":
            shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = self._hex_to_rgb(color)
            shape.line.fill.background()
        elif shape_type == "cross":
            self._add_cross_shape(slide, x, y, width, height, color)
        elif shape_type == "polygon":
            sides = element.get("pattern_data", {}).get("sides", 6)
            if sides == 6:
                shape = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, x, y, width, height)
            else:
                shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = self._hex_to_rgb(color)
            shape.line.fill.background()
        elif shape_type == "line":
            line = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, x, y, x + width, y + height
            )
            line.line.color.rgb = self._hex_to_rgb(color)
            line.line.width = Pt(2)

    def _add_cross_shape(self, slide, x, y, width, height, color: str) -> None:
        """Add a medical cross shape (two overlapping rectangles)."""
        # Vertical bar
        v_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x + width * 0.35, y,
            width * 0.3, height
        )
        v_bar.fill.solid()
        v_bar.fill.fore_color.rgb = self._hex_to_rgb(color)
        v_bar.line.fill.background()

        # Horizontal bar
        h_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x, y + height * 0.35,
            width, height * 0.3
        )
        h_bar.fill.solid()
        h_bar.fill.fore_color.rgb = self._hex_to_rgb(color)
        h_bar.line.fill.background()

    def _add_watermark(self, slide, watermark_config: Dict[str, Any]) -> None:
        """Add watermark to the slide."""
        watermark_type = watermark_config.get("type", "text")
        content = watermark_config.get("content", "")
        position = watermark_config.get("position", {"x": 10, "y": 6})

        if watermark_type == "text":
            watermark_box = slide.shapes.add_textbox(
                Inches(position["x"]), Inches(position["y"]),
                Inches(3), Inches(1)
            )
            frame = watermark_box.text_frame
            para = frame.paragraphs[0]
            para.text = content
            if para.runs:
                run = para.runs[0]
                run.font.size = Pt(24)
                run.font.bold = True
                run.font.color.rgb = RGBColor(200, 200, 200)

    @staticmethod
    def _hex_to_rgb(hex_color: str) -> RGBColor:
        """Convert hex color to RGBColor."""
        hex_color = hex_color.lstrip('#')
        return RGBColor.from_string(hex_color)


class GeometricBackgroundRenderer(BaseBackgroundRenderer):
    """Renderer for geometric patterns."""

    def supports_background_type(self, bg_type: str) -> bool:
        """Check if this renderer supports the background type."""
        return bg_type in [
            "diagonal_lines", "hexagon_grid", "geometric_pattern",
            "corporate_bars"
        ]

    def _apply_patterns(self, slide, background_design) -> None:
        """Apply geometric patterns."""
        bg_type = self._normalize_bg_type(getattr(background_design, 'type', BackgroundType.SOLID))

        if bg_type == BackgroundType.DIAGONAL_LINES:
            self._create_diagonal_lines(slide, background_design)
        elif bg_type == BackgroundType.HEXAGON_GRID:
            self._create_hexagon_grid(slide, background_design)
        elif bg_type == BackgroundType.CORPORATE_BARS:
            self._create_corporate_bars(slide, background_design)

    def _create_diagonal_lines(self, slide, background_design) -> None:
        """Create diagonal line pattern."""
        secondary_color = getattr(background_design, 'secondary_color', '#E0E0E0')
        for i in range(0, 20, 3):
            line = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(i * 0.7), Inches(0),
                Inches((i + 10) * 0.7), Inches(7.5)
            )
            line.line.color.rgb = self._hex_to_rgb(secondary_color)
            line.line.width = Pt(1)

    def _create_hexagon_grid(self, slide, background_design) -> None:
        """Create hexagon grid pattern."""
        secondary_color = getattr(background_design, 'secondary_color', '#333333')
        hex_size = 0.8
        rows, cols = 5, 8

        for row in range(rows):
            for col in range(cols):
                x = Inches(col * 1.6 + (row % 2) * 0.8)
                y = Inches(row * 1.4)
                hex_shape = slide.shapes.add_shape(
                    MSO_SHAPE.HEXAGON, x, y, Inches(hex_size), Inches(hex_size)
                )
                hex_shape.fill.background()
                hex_shape.line.color.rgb = self._hex_to_rgb(secondary_color)
                hex_shape.line.width = Pt(0.5)

    def _create_corporate_bars(self, slide, background_design) -> None:
        """Create corporate sidebar pattern."""
        secondary_color = getattr(background_design, 'secondary_color', '#E0E0E0')
        sidebar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(1), Inches(7.5)
        )
        sidebar.fill.solid()
        sidebar.fill.fore_color.rgb = self._hex_to_rgb(secondary_color)
        sidebar.line.fill.background()


class OrganicBackgroundRenderer(BaseBackgroundRenderer):
    """Renderer for organic/flowing patterns."""

    def supports_background_type(self, bg_type: str) -> bool:
        """Check if this renderer supports the background type."""
        return bg_type in ["medical_curves", "wave_design"]

    def _apply_patterns(self, slide, background_design) -> None:
        """Apply organic patterns."""
        bg_type = self._normalize_bg_type(getattr(background_design, 'type', BackgroundType.SOLID))

        if bg_type == BackgroundType.MEDICAL_CURVES:
            self._create_medical_curves(slide, background_design)

    def _create_medical_curves(self, slide, background_design) -> None:
        """Create flowing curves pattern."""
        secondary_color = getattr(background_design, 'secondary_color', '#F0FFF4')
        for i in range(3):
            curve = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(i * 4), Inches(i * 2.5),
                Inches(6), Inches(3)
            )
            curve.fill.solid()
            curve.fill.fore_color.rgb = self._hex_to_rgb(secondary_color)
            curve.line.fill.background()


class TechBackgroundRenderer(BaseBackgroundRenderer):
    """Renderer for tech/digital patterns."""

    def supports_background_type(self, bg_type: str) -> bool:
        """Check if this renderer supports the background type."""
        return bg_type in ["tech_circuit", "digital_grid"]

    def _apply_patterns(self, slide, background_design) -> None:
        """Apply tech patterns."""
        bg_type = self._normalize_bg_type(getattr(background_design, 'type', BackgroundType.SOLID))

        if bg_type == BackgroundType.TECH_CIRCUIT:
            self._create_circuit_pattern(slide, background_design)

    def _create_circuit_pattern(self, slide, background_design) -> None:
        """Create circuit board pattern."""
        secondary_color = getattr(background_design, 'secondary_color', '#2D3748')
        circuit_points = [
            (1, 1), (3, 1), (3, 3), (6, 3), (6, 5), (9, 5),
            (9, 2), (11, 2), (11, 6), (13, 6)
        ]

        for i in range(len(circuit_points) - 1):
            x1, y1 = circuit_points[i]
            x2, y2 = circuit_points[i + 1]
            line = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(x1), Inches(y1), Inches(x2), Inches(y2)
            )
            line.line.color.rgb = self._hex_to_rgb(secondary_color)
            line.line.width = Pt(2)


class BackgroundRendererFactory:
    """Factory for creating background renderers."""

    @staticmethod
    def create_renderer(bg_type: str) -> IBackgroundRenderer:
        """Create appropriate renderer for the background type."""
        if bg_type in ["diagonal_lines", "hexagon_grid", "corporate_bars"]:
            return GeometricBackgroundRenderer()
        elif bg_type in ["medical_curves", "wave_design"]:
            return OrganicBackgroundRenderer()
        elif bg_type in ["tech_circuit", "digital_grid"]:
            return TechBackgroundRenderer()
        else:
            return GeometricBackgroundRenderer()  # Default


class CompositeBackgroundRenderer(IBackgroundRenderer):
    """Composite renderer that delegates to appropriate sub-renderers."""

    def __init__(self):
        self.renderers = [
            GeometricBackgroundRenderer(),
            OrganicBackgroundRenderer(),
            TechBackgroundRenderer()
        ]

    def supports_background_type(self, bg_type: str) -> bool:
        """Check if any sub-renderer supports this type."""
        return any(renderer.supports_background_type(bg_type) for renderer in self.renderers)

    def render_background(self, slide, design_config: Dict[str, Any]) -> None:
        """Find appropriate renderer and delegate."""
        background_design = design_config.get('background_design')
        if not background_design:
            return

        bg_type = getattr(background_design, 'type', BackgroundType.SOLID)
        if hasattr(bg_type, 'value'):
            bg_type = bg_type.value

        for renderer in self.renderers:
            if renderer.supports_background_type(bg_type):
                renderer.render_background(slide, design_config)
                break